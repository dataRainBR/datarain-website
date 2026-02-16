#!/usr/bin/env python3
"""
Gera PPTs de todos os 24 cases publicados a partir do template 'Case PQ.pptx'.
Cada case gera um arquivo 'Case <Cliente>.pptx'.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import os
import sys


def replace_shape_text(shape, new_text):
    """Substitui o texto de um shape mantendo a formatação do primeiro run."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    new_paragraphs = new_text.split('\n')
    first_para = tf.paragraphs[0]
    if not first_para.runs:
        return
    ref_run = first_para.runs[0]
    ref_font_bold = ref_run.font.bold
    ref_font_size = ref_run.font.size
    ref_font_color = ref_run.font.color.rgb if ref_run.font.color and ref_run.font.color.rgb else None
    ref_font_name = ref_run.font.name
    ref_para_alignment = first_para.alignment
    ref_para_space_before = first_para.space_before
    ref_para_space_after = first_para.space_after
    ref_para_line_spacing = first_para.line_spacing
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    first_p = tf.paragraphs[0]
    for r in first_p._p.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}r'):
        first_p._p.remove(r)
    for idx, para_text in enumerate(new_paragraphs):
        if idx == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = ref_para_alignment
        if ref_para_space_before is not None:
            para.space_before = ref_para_space_before
        if ref_para_space_after is not None:
            para.space_after = ref_para_space_after
        if ref_para_line_spacing is not None:
            para.line_spacing = ref_para_line_spacing
        run = para.add_run()
        run.text = para_text
        run.font.bold = ref_font_bold
        run.font.size = ref_font_size
        if ref_font_color:
            run.font.color.rgb = ref_font_color
        if ref_font_name:
            run.font.name = ref_font_name


def replace_multiformat_text(shape, sections):
    """
    Substitui texto com múltiplas formatações.
    sections = [(text, is_title), ...]
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    title_fmt = {}
    body_fmt = {}
    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.bold and run.font.size and run.font.size >= 170000:
                if not title_fmt:
                    title_fmt = {
                        'bold': run.font.bold, 'size': run.font.size,
                        'color': run.font.color.rgb if run.font.color and run.font.color.rgb else None,
                        'name': run.font.name, 'alignment': para.alignment,
                    }
            else:
                if not body_fmt:
                    body_fmt = {
                        'bold': run.font.bold, 'size': run.font.size,
                        'color': run.font.color.rgb if run.font.color and run.font.color.rgb else None,
                        'name': run.font.name, 'alignment': para.alignment,
                    }
    if not title_fmt:
        title_fmt = body_fmt
    if not body_fmt:
        body_fmt = title_fmt
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    first_p = tf.paragraphs[0]
    for r in first_p._p.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}r'):
        first_p._p.remove(r)
    first = True
    for text, is_title in sections:
        lines = text.split('\n')
        for line in lines:
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            fmt = title_fmt if is_title else body_fmt
            para.alignment = fmt.get('alignment')
            run = para.add_run()
            run.text = line
            run.font.bold = fmt.get('bold')
            run.font.size = fmt.get('size')
            if fmt.get('color'):
                run.font.color.rgb = fmt.get('color')
            if fmt.get('name'):
                run.font.name = fmt.get('name')


def replace_image(slide, shape_id, new_image_path):
    """Substitui a imagem de um shape mantendo posição e tamanho. Converte WEBP para PNG se necessário."""
    actual_path = new_image_path
    if new_image_path.lower().endswith('.webp'):
        from PIL import Image
        png_path = new_image_path.rsplit('.', 1)[0] + '_converted.png'
        if not os.path.exists(png_path):
            img = Image.open(new_image_path)
            img.save(png_path, 'PNG')
        actual_path = png_path

    for shape in slide.shapes:
        if shape.shape_id == shape_id and shape.shape_type == 13:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(actual_path, left, top, width, height)
            return True
    return False


# ============================================================
# DADOS DE TODOS OS 24 CASES
# ============================================================
# Cada case é um dict com os campos do template PPT:
#   headline, titulo_secao1, texto_cliente, titulo_secao2, texto_desafio,
#   titulo_secao3, texto_solucao, texto_voce_sabia, texto_transicao,
#   titulo_arquitetura, texto_arquitetura, titulo_resultados, texto_resultados,
#   titulo_depoimento (ou None), texto_depoimento, texto_autor_depoimento,
#   titulo_continuacao, texto_continuacao,
#   logo_path, output_name

CASES = [
    # ---- CASE 1: PQ Corp (template original, pular) ----

    # ---- CASE 2: Bitz (migração) ----
    {
        "headline": "Bitz acelera sua transformação digital com migração para nuvem AWS em apenas 60 dias.",
        "titulo_secao1": "COMO O BITZ MIGROU PARA A AWS EM TEMPO RECORDE",
        "texto_cliente": "O Bitz, carteira digital do Banco Bradesco, chegou para facilitar e organizar a vida financeira de seus clientes. Por meio de um aplicativo, é possível guardar dinheiro e cartões, realizar transações, receber pagamentos, transferir dinheiro e aproveitar descontos exclusivos. Prático, o Bitz oferece pagamento via QRCode nas máquinas da Cielo, com mais de 1,6 milhão de estabelecimentos cadastrados em todo o país.",
        "titulo_secao2": "QUANDO O CRESCIMENTO EXIGE VELOCIDADE",
        "texto_desafio": "O Bitz precisava de volume de dados para suportar suas transações, especialmente por conta da Black Friday. A melhor solução apontou para o caminho de se trabalhar com autonomia de operação adotando ambientes na nuvem AWS.\n\nA migração precisava ser realizada com agilidade, de forma assertiva e precisa. Os fatores críticos eram agilidade para colocar o novo ambiente no ar, elasticidade para acompanhar o rápido crescimento do número de clientes, e fazê-los com muita segurança e investimentos razoáveis.",
        "titulo_secao3": "PARCERIA DATARAIN: MIGRAÇÃO EM 60 DIAS",
        "texto_solucao": "A estratégia apostou em três pilares:\nImplementação de Landing Zone com segurança e conformidade\nMigração completa em apenas 60 dias com metodologia validada\nTestes de carga e otimização de custos para suportar o crescimento acelerado",
        "texto_voce_sabia": "Segundo a ACI Worldwide, o Brasil é o quarto país do mundo em pagamentos em tempo real, com 8,7 bilhões de transações em 2021. A expectativa é que esse tipo de transação chegue a 82,4 bilhões ao ano até 2026, gerando US$ 37,6 bilhões adicionais de produção econômica.",
        "texto_transicao": "Com o desafio de migrar toda a infraestrutura para a nuvem em tempo recorde, a proposta era clara: garantir elasticidade, segurança e alta disponibilidade para suportar o crescimento acelerado da base de clientes.",
        "titulo_arquitetura": "MIGRAÇÃO COMPLETA: DO LEGADO À NUVEM AWS",
        "texto_arquitetura": "A migração incluiu assessment do ambiente legado, configuração de Landing Zone, distribuição de cargas de trabalho e gerenciamento de identidades. No Dia D, foram migradas bases MS SQL e DynamoDB, virada de DNS, testes de liberação em Android e iOS, e setup de CI/CD.\n\nA infraestrutura utiliza Amazon S3, AWS Lambda, Amazon RDS, EC2, CloudWatch, CloudFront, VPC, DMS, ELB, SES, ECS, DynamoDB, Route 53, CloudTrail e SQS.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: CRESCIMENTO DE 150%",
        "texto_resultados": "Desde a migração, o Bitz alcançou 3,4 milhões de downloads/mês e 8 milhões de logins mensais. A base de clientes cresceu 150%. O pós-migração não trouxe qualquer impacto negativo aos usuários.\n\nA migração em apenas 60 dias trouxe elasticidade, agilidade e maior segurança para acompanhar o crescimento da base de clientes e implantar novos serviços.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "RUMO AO TÍTULO DE UNICÓRNIO",
        "texto_continuacao": "Como um dos principais players do setor, o Bitz já almeja o título de Unicórnio. Trabalha ativamente para acompanhar o crescimento acelerado da categoria de contas de pagamento, elevando seu patamar de segurança e trazendo facilidade em um aplicativo grátis e completo.\n\nArmazenamento, disponibilidade, segurança, elasticidade e alta performance são vantagens angariadas com a adesão à AWS, fundamentais para o funcionamento do sistema.",
        "logo_path": "public/content-images/cases/bitz-acelera-sua-transformacao-digital-com-migracao-para-nuvem-aws-em-apenas-60-dias/5.png.webp",
        "output_name": "Case Bitz Migracao",
    },
    # ---- CASE 3: Bitz (DevOps) ----
    {
        "headline": "DevOps garante sucesso na migração do Bitz, carteira digital do Bradesco.",
        "titulo_secao1": "COMO O DEVOPS GARANTIU A MIGRAÇÃO DO BITZ",
        "texto_cliente": "O Bitz é a carteira digital do Banco Bradesco, uma inovação da instituição financeira que conta com quase 80 anos de mercado. A plataforma moderna de conta de pagamento grátis é operada via aplicativo de smartphone, buscando oferecer uma experiência contemporânea e segura para seus usuários, refletindo os valores de confiabilidade e credibilidade do Bradesco.",
        "titulo_secao2": "QUANDO A MIGRAÇÃO EXIGE CONTINUIDADE OPERACIONAL",
        "texto_desafio": "O desafio era migrar a infraestrutura do Bitz para a nuvem da AWS em curto prazo, visando transacionar em um novo ambiente durante a Black Friday. Era necessário garantir a continuidade das operações financeiras de forma inovadora, com menor custo e sem abrir mão da segurança.\n\nA abordagem DevOps foi fundamental para garantir o sucesso da implantação e migração para a nuvem AWS.",
        "titulo_secao3": "PARCERIA DATARAIN: DEVOPS EM 60 DIAS",
        "texto_solucao": "A estratégia apostou em três pilares:\nAWS ECS com Auto Scaling para dimensionar containers automaticamente\nBalanceador de carga integrado ao ECS para distribuir tráfego\nBanco de dados AWS RDS com Multi-AZ para alta disponibilidade",
        "texto_voce_sabia": "A cultura DevOps integra equipes de desenvolvimento e operações de TI, visando a entrega contínua de valor. Com Auto Scaling, se um container falhar, outro é iniciado automaticamente para substituí-lo, garantindo disponibilidade contínua do serviço.",
        "texto_transicao": "Com o desafio de garantir continuidade operacional durante a migração, a proposta era clara: implementar uma infraestrutura DevOps escalável e altamente disponível em apenas 60 dias.",
        "titulo_arquitetura": "INFRAESTRUTURA DEVOPS ESCALÁVEL NA AWS",
        "texto_arquitetura": "A solução utilizou AWS ECS com Auto Scaling para dimensionar containers com base em métricas de CPU e tráfego. O gerenciamento de tarefas monitora todos os containers, reiniciando automaticamente em caso de falha.\n\nO frontend é armazenado no AWS S3 e entregue pelo CloudFront, ambos escaláveis e altamente disponíveis. Os bancos de dados críticos são configurados com Multi-AZ no AWS RDS, garantindo failover automático.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: DEVOPS QUE ENTREGA",
        "texto_resultados": "Dimensionamento automático dos containers com base em métricas específicas. Monitoramento centralizado com capacidade de reiniciar containers em caso de falhas. Distribuição eficiente do tráfego entre containers.\n\nArmazenamento seguro de dados e logs fora dos containers. Escalabilidade e alta disponibilidade do frontend. Configuração resiliente do banco de dados com failover automático.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "ESCALABILIDADE E RESILIÊNCIA",
        "texto_continuacao": "A implementação permitiu uma infraestrutura DevOps escalável e altamente disponível. A cultura DevOps promoveu colaboração entre as equipes, permitindo que o Bitz atingisse seus objetivos de negócio.\n\nA nova infraestrutura permitiu que o Bitz colocasse seus serviços no ar antes do previsto, validando conceitos técnicos e preparando o terreno para a expansão da base de clientes.",
        "logo_path": "public/content-images/cases/devops-garante-sucesso-na-migracao-do-bitz-carteira-digital-do-bradesco/5.png.webp",
        "output_name": "Case Bitz DevOps",
    },
    # ---- CASE 4: Bitz (Observabilidade) ----
    {
        "headline": "Bitz melhora MTTI e MTTR com plataforma de observabilidade integrada à AWS.",
        "titulo_secao1": "COMO O BITZ REVOLUCIONOU SEU MONITORAMENTO",
        "texto_cliente": "O Bitz, carteira digital do Banco Bradesco, oferece pagamento via QRCode nas máquinas da Cielo, com mais de 1,6 milhão de estabelecimentos cadastrados. Com 12 milhões de clientes cadastrados, a plataforma de pagamentos precisava de monitoramento contínuo e proativo para garantir a qualidade do serviço.",
        "titulo_secao2": "QUANDO A VISIBILIDADE FAZ A DIFERENÇA",
        "texto_desafio": "O desafio era melhorar o acompanhamento das fases de desenvolvimento, testes, homologação e produção, através do monitoramento e rastreamento contínuos dos microserviços internos e de terceiros.\n\nEra necessário identificar preditivamente comportamentos anômalos, erros e ineficiências que afetavam a qualidade de serviço. Situações que antes só eram identificadas em post-mortem precisavam ser antecipadas.",
        "titulo_secao3": "PARCERIA DATARAIN: OBSERVABILIDADE INTELIGENTE",
        "texto_solucao": "A estratégia apostou em três pilares:\nSolução Best-of-The-Breed de observabilidade integrada à plataforma de pagamentos\nAlarmes proativos e dashboards técnicos e de negócios\nMonitoramento de microserviços, APIs, bancos de dados e containers",
        "texto_voce_sabia": "MTTI (Mean Time To Identify) e MTTR (Mean Time To Resolve) são métricas críticas para plataformas de pagamento. Reduzir esses tempos significa identificar e resolver incidentes antes que impactem os clientes, garantindo a continuidade do serviço.",
        "texto_transicao": "Com o desafio de monitorar uma plataforma com 12 milhões de clientes, a proposta era clara: implementar observabilidade inteligente para antecipar problemas e reduzir o impacto sobre os serviços.",
        "titulo_arquitetura": "OBSERVABILIDADE: VISÃO COMPLETA DA OPERAÇÃO",
        "texto_arquitetura": "A solução integrou uma plataforma de observabilidade Best-of-The-Breed à infraestrutura na AWS, com serviços de suporte e sustentação da dataRain. O monitoramento abrange microserviços, dependências entre APIs, bancos de dados e containers.\n\nAlarmes proativos e dashboards técnicos e de negócios permitem identificar e antecipar ações corretivas com precisão.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: ANTECIPAÇÃO E PRECISÃO",
        "texto_resultados": "Redução significativa do MTTI e MTTR na plataforma de pagamentos com 12 milhões de clientes. Correlação superior a 95% entre causa-e-efeito de incidentes e comportamentos anômalos.\n\nMinimização do impacto sobre os serviços prestados, identificando proativamente anomalias e se antecipando a eventos que antes impactavam negativamente as métricas de qualidade.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "QUALIDADE CONTÍNUA",
        "texto_continuacao": "A plataforma de observabilidade transformou a forma como o Bitz monitora e responde a incidentes. Com dashboards em tempo real e alarmes proativos, a equipe consegue antecipar problemas antes que afetem os clientes.\n\nA correlação precisa entre causa e efeito permite ações corretivas direcionadas, elevando continuamente a qualidade do serviço.",
        "logo_path": "public/content-images/cases/melhorando-o-mtti-e-mttr-atraves-de-uma-plataforma-de-observabilidade/5.png.webp",
        "output_name": "Case Bitz Observabilidade",
    },
    # ---- CASE 5: Prefeitura de Americana ----
    {
        "headline": "Serviços digitais diminuem em 50% a espera no atendimento público em Americana.",
        "titulo_secao1": "COMO AMERICANA DIGITALIZOU SEUS SERVIÇOS PÚBLICOS",
        "texto_cliente": "A prefeitura de Americana, localizada no estado de São Paulo, buscava uma informatização completa dos serviços prestados por todos os órgãos e secretarias da administração municipal. A plataforma 1Doc foi a solução escolhida para transformar digitalmente o atendimento ao cidadão.",
        "titulo_secao2": "QUANDO O PAPEL ATRASA O ATENDIMENTO",
        "texto_desafio": "Até março de 2020, a prefeitura tramitava todos os processos de forma analógica, gerando mais gastos aos cofres públicos, tornando a resolução de demandas mais burocrática, agredindo o meio ambiente e resultando em tempo de espera alto na tratativa das solicitações da população.\n\nA necessidade de digitalização se tornou urgente com a pandemia, exigindo atendimento remoto e eficiente.",
        "titulo_secao3": "PARCERIA DATARAIN: TRANSFORMAÇÃO DIGITAL COMPLETA",
        "texto_solucao": "A estratégia apostou em três pilares:\nPlataforma 1Doc com arquitetura em alta disponibilidade na AWS\nPipeline de desenvolvimento com CI/CD usando serviços AWS\nImplantação 100% remota com treinamento à distância",
        "texto_voce_sabia": "A plataforma 1Doc utiliza Amazon Aurora, Amazon S3, ElastiCache para Redis, ECS para containers escaláveis, e um pipeline completo de CI/CD com CodePipeline, CodeBuild e CodeDeploy. A migração do RDS para Aurora trouxe melhor performance sem necessidade de alocar disco para IOPS.",
        "texto_transicao": "Com o desafio de digitalizar completamente os serviços municipais, a proposta era clara: eliminar o papel, reduzir custos e acelerar o atendimento ao cidadão com tecnologia em nuvem.",
        "titulo_arquitetura": "PLATAFORMA DIGITAL: DO PAPEL À NUVEM",
        "texto_arquitetura": "A arquitetura utiliza Amazon Aurora para dados, S3 para anexos e código estático, ElastiCache para Redis como cache, ECS para containers escaláveis, SQS para operações assíncronas e CloudWatch para monitoramento.\n\nO pipeline de CI/CD integra Bitbucket, CodePipeline, CodeBuild, ECR e CodeDeploy, com staging em EC2 e produção em ECS.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: ATENDIMENTO DIGITAL EM ESCALA",
        "texto_resultados": "Mais de 450.000 atendimentos digitais (85% do total). Redução de 10x no tempo de emissão de certidão negativa. 99,9% dos projetos de construção analisados digitalmente. 50% de redução no tempo médio de aprovação de projetos.\n\n630 serviços online disponíveis. Mais de 10 milhões de impressões poupadas, resultando em mais de R$1.000.000 em economia aos cofres públicos.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "GOVERNO DIGITAL EM EXPANSÃO",
        "texto_continuacao": "A transformação digital permitiu ao município adotar comunicação digital, possibilitando ao cidadão contato e pedido de demandas a qualquer hora do dia, sem necessidade de deslocamento físico.\n\nA implantação aconteceu de forma 100% remota, com reuniões por videoconferência e treinamento via educação à distância, demonstrando que a modernização do setor público é possível e acessível.",
        "logo_path": "public/content-images/cases/cases-prefeitura-americana/logo-pma-eleicao-2.png",
        "output_name": "Case Prefeitura Americana",
    },
    # ---- CASE 6: INPE Biomassa ----
    {
        "headline": "INPE reduz de 3.000 para 40 horas o processamento do mapa de biomassa amazônica.",
        "titulo_secao1": "COMO A NUVEM ACELEROU A CIÊNCIA AMAZÔNICA",
        "texto_cliente": "O Instituto Nacional de Pesquisas Espaciais (INPE) é um dos principais centros de pesquisa do Brasil, reconhecido mundialmente por sua atuação no monitoramento ambiental e climático. O projeto EBA (Estimativa de Biomassa da Amazônia) mapeou a biomassa do Bioma Amazônico combinando dados LiDAR de mais de 1.000 voos com informações de satélites.",
        "titulo_secao2": "QUANDO A CIÊNCIA ESBARRA NO PROCESSAMENTO",
        "texto_desafio": "Cada voo gerava 11 GB com cerca de 6,5 bilhões de registros. Com mais de 1.000 voos, isso resultou em cerca de 7 trilhões de registros. Para calcular a incerteza de cada pixel, era necessário gerar mil versões diferentes dos mapas de biomassa.\n\nMesmo com bons servidores locais, processar todo esse volume levaria mais de 3.000 horas, inviabilizando o projeto dentro do prazo.",
        "titulo_secao3": "PARCERIA DATARAIN: HPC NA NUVEM AWS",
        "texto_solucao": "A estratégia apostou em três pilares:\nAmazon EC2 HPC com 4 instâncias de 64 vCPUs e 256 GB RAM cada\nAlgoritmo Random Forest com plataforma H2O e Python\nArmazenamento em S3 com sistema de arquivos compartilhado EFS",
        "texto_voce_sabia": "O mapa de incertezas permite determinar onde o modelo foi menos preciso, subsidiando estudos complementares. Com ele, a atualização futura do mapa pode ser feita com apenas 15% dos voos iniciais, que custaram cerca de R$3.000.000 à época.",
        "texto_transicao": "Com o desafio de processar 7 trilhões de registros para gerar mil mapas de biomassa, a proposta era clara: usar computação de alta performance na nuvem para viabilizar a ciência em escala.",
        "titulo_arquitetura": "HPC NA NUVEM: PROCESSAMENTO EM ESCALA",
        "texto_arquitetura": "A arquitetura combinou Amazon EC2 HPC (4 instâncias com 64 vCPUs e 256 GB RAM), Amazon S3 para armazenamento, Amazon EBS para volumes de disco, Amazon EFS para sistema de arquivos compartilhado, Route 53 para DNS e VPC para rede isolada.\n\nA modelagem estatística utilizou o algoritmo Random Forest com a plataforma H2O e Python para calcular a incerteza de cada pixel do mapa.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: DE 3.000 PARA 40 HORAS",
        "texto_resultados": "O tempo de processamento caiu de 2.000-3.000 horas para apenas 40 horas com a infraestrutura em nuvem. O INPE consolidou todo o trabalho em seu servidor local em apenas dois dias.\n\nO mapa de incertezas permite atualização futura com apenas 15% dos voos iniciais, reduzindo drasticamente os custos da etapa mais cara do projeto.",
        "texto_depoimento": '"Para estimar a incerteza do nosso conhecimento sobre a biomassa da Amazônia, eu precisava gerar mil mapas. Com a estrutura que tínhamos, isso levaria de 2 a 3 mil horas. Com o apoio da AWS e da dataRain, conseguimos processar tudo em apenas 40 horas."',
        "texto_autor_depoimento": "- Mauro Assis, pesquisador do CCST/INPE",
        "titulo_continuacao": "LEGADO CIENTÍFICO COM IMPACTO ESTRATÉGICO",
        "texto_continuacao": "Concluído em 2019, o projeto deixou um legado técnico e científico: a floresta foi mapeada com mais precisão e responsabilidade. Os dados gerados seguem sendo referência para pesquisas, inventários de carbono e políticas públicas.\n\nUm exemplo claro de como nuvem, ciência e colaboração transformam conhecimento em impacto real para discussões internacionais sobre clima.",
        "logo_path": "public/content-images/cases/como-a-nuvem-acelerou-o-mapa-de-incertezas-da-biomassa-amazonica/unnamed.png",
        "output_name": "Case INPE Biomassa",
    },
    # ---- CASE 7: BluePay ----
    {
        "headline": "BluePay constrói plataforma de pagamentos PIX escalável e segura na AWS.",
        "titulo_secao1": "COMO A BLUEPAY ESCALOU SUA PLATAFORMA PIX",
        "texto_cliente": "A BluePay Fintech é uma plataforma de pagamentos especializada em PIX, focada em escala, alta disponibilidade, segurança e custos adequados. A empresa reconheceu a importância da bancarização de profissionais autônomos e prestadores de serviços, visando oferecer acesso simplificado a serviços financeiros.",
        "titulo_secao2": "QUANDO A VELOCIDADE DO MERCADO EXIGE AGILIDADE",
        "texto_desafio": "O principal desafio era reduzir o tempo de desenvolvimento e implementação de novos recursos na plataforma, sem comprometer segurança ou estabilidade. A BluePay precisava responder rapidamente às demandas do mercado de pagamentos.\n\nEra necessário escalar recursos de forma flexível, garantir alta disponibilidade e manter controles rígidos de segurança e conformidade.",
        "titulo_secao3": "PARCERIA DATARAIN: DEVOPS + SERVERLESS",
        "texto_solucao": "A estratégia apostou em três pilares:\nArquitetura serverless com AWS Lambda, API Gateway e DynamoDB/RDS\nCI/CD automatizado com GitHub Pipeline e Serverless Framework\nSegurança com princípio do menor privilégio e criptografia KMS",
        "texto_voce_sabia": "A arquitetura serverless permite que a BluePay pague apenas pelo processamento efetivamente utilizado, sem necessidade de provisionar servidores. Com AWS Lambda, o dimensionamento é automático, acompanhando picos de transações PIX em tempo real.",
        "texto_transicao": "Com o desafio de construir uma plataforma de pagamentos PIX que acompanhasse o crescimento acelerado, a proposta era clara: combinar DevOps com arquitetura serverless para máxima agilidade e segurança.",
        "titulo_arquitetura": "SERVERLESS: ESCALABILIDADE SOB DEMANDA",
        "texto_arquitetura": "A arquitetura utiliza AWS Lambda para processamento, API Gateway para acessibilidade, DynamoDB/RDS como camada de dados, CloudWatch para monitoramento, CloudFront para distribuição do frontend com HTTPS, e KMS para criptografia.\n\nSegurança com Identity Center (SSO), VPC Flow Logs, criptografia padrão em RDS e DynamoDB, e tráfego protegido com TLS/SSL via ACM. RPO de 24h e RTO de 12h.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: AGILIDADE E SEGURANÇA",
        "texto_resultados": "CI/CD possibilitou lançamentos frequentes de funcionalidades, mantendo a BluePay na vanguarda da inovação em pagamentos. Automação de testes e detecção precoce de erros resultaram em experiência mais estável.\n\nArquitetura serverless permite dimensionar a infraestrutura automaticamente em picos de tráfego. Cultura DevOps integrou desenvolvimento, operações e segurança com comunicação contínua.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "INOVAÇÃO CONTÍNUA EM PAGAMENTOS",
        "texto_continuacao": "A combinação de DevOps e serverless na AWS permitiu à BluePay criar um ambiente de pagamentos escalável, seguro e confiável. O foco em automação e monitoramento garante a entrega de serviços de alta qualidade.\n\nTestes automatizados de segurança, monitoramento contínuo e controles de acesso granulares protegem os dados dos clientes e garantem conformidade.",
        "logo_path": "public/content-images/cases/devops-com-mais-seguranca-e-escalabilidade-para-bluepay/BluePay.webp",
        "output_name": "Case BluePay",
    },
    # ---- CASE 8: Banco Ótimo ----
    {
        "headline": "Banco Ótimo constrói plataforma financeira 100% online com DevOps e serverless na AWS.",
        "titulo_secao1": "COMO O BANCO ÓTIMO NASCEU NA NUVEM",
        "texto_cliente": "O Banco Ótimo (Ótimo SCD S/A) é uma plataforma financeira 100% online que oferece serviços simplificados e acessíveis. Autorizado pelo Banco Central, conta com a custódia e experiência da CODEPE CVC S/A, empresa com 34 anos de atuação no mercado financeiro e de capitais.",
        "titulo_secao2": "QUANDO A INCLUSÃO FINANCEIRA EXIGE TECNOLOGIA",
        "texto_desafio": "O Banco Ótimo reconheceu a importância da bancarização de profissionais autônomos e prestadores de serviços. Precisava agilizar o desenvolvimento e implementação de soluções financeiras personalizadas para esse público.\n\nEra necessária uma abordagem que permitisse responder rapidamente às demandas do mercado, com escalabilidade, segurança e custos otimizados.",
        "titulo_secao3": "PARCERIA DATARAIN: BANCO DIGITAL NA AWS",
        "texto_solucao": "A estratégia apostou em três pilares:\nArquitetura serverless com AWS Lambda, API Gateway e DynamoDB/RDS\nIdentity Center (SSO) com credenciais temporárias e menor privilégio\nServerless Framework com CloudFormation para infraestrutura como código",
        "texto_voce_sabia": "O Banco Ótimo opera com RPO de até 24 horas (snapshots diários) e RTO de até 12 horas, garantindo recuperação rápida em caso de falhas. A criptografia padrão é habilitada em todos os bancos de dados, com chaves armazenadas no AWS KMS.",
        "texto_transicao": "Com o desafio de criar um banco digital 100% online para inclusão financeira, a proposta era clara: construir uma plataforma serverless escalável, segura e com custos otimizados na AWS.",
        "titulo_arquitetura": "BANCO DIGITAL: SERVERLESS E SEGURO",
        "texto_arquitetura": "A arquitetura utiliza AWS Lambda para processamento, API Gateway para acessibilidade, DynamoDB/RDS como camada de dados. CloudWatch com painel personalizado para monitoramento de infraestrutura, aplicação e erros.\n\nSegurança com Identity Center (SSO), VPC Flow Logs em produção, criptografia TLS/SSL via ACM, CloudFront para HTTPS, e criptografia padrão em RDS e DynamoDB com KMS.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: BANCO ESCALÁVEL E SEGURO",
        "texto_resultados": "Agilidade e velocidade com CI/CD para lançamento frequente de funcionalidades. Confiabilidade e estabilidade com automação de processos e testes automatizados. Escalabilidade e flexibilidade com arquitetura serverless.\n\nColaboração entre equipes de desenvolvimento, operações e segurança. Segurança reforçada com testes automatizados e controles de acesso granulares.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "INCLUSÃO FINANCEIRA DIGITAL",
        "texto_continuacao": "A adoção de DevOps e serverless na AWS permitiu ao Banco Ótimo criar um ambiente bancário escalável, seguro e confiável. O foco em automação e monitoramento garante a entrega de serviços de alta qualidade.\n\nA parceria com a dataRain possibilitou uma implantação bem-sucedida, com menor tempo de entrega, maior agilidade e confiabilidade para os clientes.",
        "logo_path": "public/content-images/cases/devops-traz-mais-seguranca-e-escalabilidade-ao-banco-otimo/4.png.webp",
        "output_name": "Case Banco Otimo",
    },
    # ---- CASE 9: Pakman ----
    {
        "headline": "Pakman constrói infraestrutura serverless escalável para logística de last mile.",
        "titulo_secao1": "COMO A PAKMAN ESCALOU SUA LOGÍSTICA",
        "texto_cliente": "A Pakman é uma Logtech especializada em serviços de Last Mile. Desde desenvolvimento à execução de soluções para empresas que possuem necessidade de entrega com excelência, pontualidade e rastreamento para otimizar a experiência do cliente no e-commerce.",
        "titulo_secao2": "QUANDO O E-COMMERCE CRESCE EXPONENCIALMENTE",
        "texto_desafio": "O maior desafio era a necessidade de entrega eficiente e pontual. Antes da AWS, a Pakman enfrentava gestão inadequada do tempo de entrega, falta de visibilidade em tempo real e limitações na otimização dinâmica das rotas.\n\nA incapacidade de oferecer entregas rápidas e precisas resultava em insatisfação dos clientes, que buscam cada vez mais conveniência e agilidade nas últimas etapas do processo logístico.",
        "titulo_secao3": "PARCERIA DATARAIN: LOGÍSTICA SERVERLESS",
        "texto_solucao": "A estratégia apostou em três pilares:\nArquitetura serverless com Lambda, Step Functions e API Gateway\nDynamoDB Stream e SQS Lambda Trigger para processamento em tempo real\nFargate e EKS para containers com escalabilidade automática",
        "texto_voce_sabia": "A Pakman utiliza Lambda, Step Functions, SQS, CDK, RDS Aurora MySQL, DynamoDB, S3, API Gateway, SNS, IAM, CloudFront, Route 53, Fargate e CloudFormation. Essa combinação permite escalar automaticamente conforme o volume de pedidos do e-commerce.",
        "texto_transicao": "Com o desafio de acompanhar o crescimento exponencial dos pedidos de logística do e-commerce, a proposta era clara: construir uma arquitetura robusta e altamente escalável na AWS.",
        "titulo_arquitetura": "LOGÍSTICA EM ESCALA: SERVERLESS NA AWS",
        "texto_arquitetura": "A arquitetura utiliza AWS Lambda para processamento, API Gateway para acessibilidade, DynamoDB/RDS como camada de dados. Step Functions orquestram fluxos complexos, SQS gerencia filas de mensagens e DynamoDB Stream processa eventos em tempo real.\n\nSegurança com VPC Flow Logs, criptografia TLS/SSL via ACM, CloudFront para HTTPS, e criptografia padrão em RDS e DynamoDB com KMS. RPO de 24h e RTO de 12h.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: LOGÍSTICA ÁGIL E CONFIÁVEL",
        "texto_resultados": "Infraestrutura que cresce de maneira sinérgica com o aumento das operações do e-commerce. Experiência de logística ágil e confiável para os clientes finais.\n\nEscalabilidade automática para lidar com picos de demanda. Visibilidade em tempo real das entregas e otimização dinâmica de rotas.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "CRESCIMENTO SEM LIMITES",
        "texto_continuacao": "A arquitetura serverless na AWS permite à Pakman escalar sua operação de forma transparente, acompanhando o crescimento exponencial do e-commerce brasileiro.\n\nA combinação de serviços gerenciados da AWS com práticas DevOps garante entregas rápidas, precisas e rastreáveis, elevando a satisfação dos clientes finais.",
        "logo_path": "public/content-images/cases/escalabilidade-com-devops-na-pakman/pakman-logo.png",
        "output_name": "Case Pakman",
    },
    # ---- CASE 10: HC Genômica ----
    {
        "headline": "HCFMUSP migra 367 TB de dados genômicos para a AWS com apoio da dataRain.",
        "titulo_secao1": "COMO O HC TRANSFORMOU SUA PESQUISA GENÔMICA",
        "texto_cliente": "O Hospital das Clínicas da Faculdade de Medicina da USP (HCFMUSP), inaugurado em 1944, é um complexo hospitalar de referência vinculado à Secretaria de Estado da Saúde e à FMUSP. No contexto das pesquisas genômicas, a UNIFESP conduziu projetos com apoio do HCFMUSP, aproximando descobertas científicas da prática clínica.",
        "titulo_secao2": "QUANDO OS DADOS GENÔMICOS EXIGEM ESCALA",
        "texto_desafio": "Hospitais e centros de pesquisa precisam armazenar e processar volumes massivos de dados sensíveis, mantendo governança e conformidade regulatória. Os principais gargalos incluíam armazenamento seguro de grandes volumes, custo e lentidão de ambientes on-premises.\n\nA complexidade de governança (incluindo LGPD) e a dificuldade de escalabilidade para acompanhar a evolução das pesquisas tornavam a modernização urgente.",
        "titulo_secao3": "PARCERIA DATARAIN: GENÔMICA NA NUVEM",
        "texto_solucao": "A estratégia apostou em três pilares:\nTransferência segura de 367 TB via AWS DataSync\nAmazon S3 como repositório central organizado por tipo de dado\nAvaliação comparativa entre EC2 e AWS HealthOmics para processamento",
        "texto_voce_sabia": "O projeto envolveu a migração de 367 TB de dados genômicos, incluindo dados previamente armazenados em Glacier. A infraestrutura na AWS permite ingestão contínua de novos dados e processamento de pipelines genômicos diretamente na nuvem.",
        "texto_transicao": "Com o desafio de modernizar o ambiente de dados genômicos do maior complexo hospitalar da América Latina, a proposta era clara: criar uma base escalável e segura para pesquisa genômica e medicina personalizada.",
        "titulo_arquitetura": "GENÔMICA EM ESCALA: DADOS NA NUVEM AWS",
        "texto_arquitetura": "A arquitetura combina Amazon S3 como repositório central organizado por tipo de dado, AWS DataSync para migração inicial e ingestão contínua, Amazon EC2 para executar pipelines genômicos na nuvem.\n\nMonitoramento e alertas de custos garantem controle e visibilidade financeira. A governança estruturada assegura conformidade com LGPD e auditoria contínua.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: BASE MODERNA PARA PESQUISA",
        "texto_resultados": "Migração precisa e sem perdas de 367 TB, com estabilidade e desempenho do novo ambiente. Dados organizados e prontos para pesquisa em escala, com ingestão contínua.\n\nRedução de custos operacionais ao otimizar processamento e armazenamento. Escalabilidade e agilidade para expandir pipelines e acelerar análises genômicas.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "MEDICINA PERSONALIZADA DO FUTURO",
        "texto_continuacao": "Os próximos passos incluem benchmarking entre processamento on-premises, EC2 e AWS HealthOmics, preparando o hospital para expandir o uso da nuvem e fortalecer a interoperabilidade de dados.\n\nAo unir a capacidade clínica do HCFMUSP, a colaboração acadêmica com a UNIFESP e a expertise da dataRain, o projeto criou condições para transformar dados genômicos em insights aplicáveis.",
        "logo_path": "public/content-images/cases/genomica-em-escala-a-transformacao-digital-no-hospital-das-clinicas/logo-hc.png",
        "output_name": "Case HC Genomica",
    },
    # ---- CASE 11: Autoglass (já gerado separadamente, incluir aqui também) ----
    {
        "headline": "Grupo Autoglass eleva governança e segurança na AWS com implementação de Landing Zone.",
        "titulo_secao1": "COMO O GRUPO AUTOGLASS ELEVOU SUA GOVERNANÇA NA AWS",
        "texto_cliente": "O Grupo Autoglass é um ecossistema de soluções que vai além do setor automotivo. Com a Autoglass, referência em vidros e peças automotivas presente em mais de 85 lojas no Brasil e na Colômbia, o grupo expandiu sua atuação com a Maxpar, especializada em assistências para seguradoras, e a B4B.tech, unidade de tecnologia focada em digitalização e BPO de operações de seguros.",
        "titulo_secao2": "QUANDO O CRESCIMENTO ESBARRA NA GOVERNANÇA",
        "texto_desafio": "O Grupo operava com cerca de seis contas AWS compartilhadas entre cinco Business Units, o que gerava dificuldades para garantir isolamento entre recursos e segurança das aplicações. A falta de controle sobre os gastos de cada BU dificultava a alocação de custos e a governança financeira.\n\nA infraestrutura apresentava controles de segurança insuficientes: ausência de controle sobre regiões AWS utilizadas, descentralização dos logs de segurança e falta de políticas preventivas e detectivas.",
        "titulo_secao3": "PARCERIA DATARAIN: IMPLEMENTANDO A LANDING ZONE",
        "texto_solucao": "A estratégia apostou em três pilares:\nSegmentar o ambiente com OUs dedicadas, garantindo isolamento e autonomia por BU\nConfigurar serviços nativos de segurança como CloudTrail, GuardDuty e Security Hub\nCapacitar o time interno com transferência de conhecimento e trabalho colaborativo",
        "texto_voce_sabia": "Segundo a AWS, empresas que implementam Landing Zones reduzem em até 50% o tempo de provisionamento de novas contas e ambientes, além de diminuir significativamente o risco de incidentes de segurança por configurações incorretas.",
        "texto_transicao": "Com o desafio de estruturar uma governança robusta para múltiplas unidades de negócio, a proposta era clara: garantir isolamento, segurança e rastreabilidade de custos com Landing Zone na AWS.",
        "titulo_arquitetura": "LANDING ZONE: ESTRUTURA SEGURA E ESCALÁVEL",
        "texto_arquitetura": "A solução incluiu OUs segmentadas: uma exclusiva para a conta Master, outra para infraestrutura e uma para contas excluídas. Contas dedicadas para backups, IAM Identity Center e configurações de rede.\n\nServiços nativos como CloudTrail, Config, GuardDuty, Security Hub e Backup elevaram a maturidade de segurança. Políticas de backups e tags centralizam cópias de segurança e rastreiam custos por BU.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: DA COMPLEXIDADE À GOVERNANÇA TOTAL",
        "texto_resultados": "BUs em contas AWS exclusivas com isolamento e segurança. SCPs e Control Tower aumentaram a governança. GuardDuty, Security Hub e CloudTrail elevaram segurança e monitoramento.\n\nPolítica de tags para alocação precisa de custos por BU. Account Factory agilizou criação de novas contas. Time interno capacitado para gerenciar a infraestrutura com autonomia.",
        "texto_depoimento": '"Um dos grandes diferenciais desse projeto foi o atendimento próximo e colaborativo da equipe da dataRain. Trabalharam lado a lado com nossos profissionais, capacitando-os e transferindo conhecimento ao longo de todo o processo."',
        "texto_autor_depoimento": "- Grupo Autoglass",
        "titulo_continuacao": "AUTONOMIA E ESCALABILIDADE",
        "texto_continuacao": "O trabalho colaborativo e a transferência de conhecimento garantiram que o time do Grupo Autoglass estivesse preparado para gerenciar e operar a nova infraestrutura com autonomia.\n\nHoje, a organização opera com um nível de segurança acima da média, com um time maduro e capacitado, reforçando a posição de liderança no setor automotivo.",
        "logo_path": "public/content-images/cases/grupo-autoglass-eleva-governanca-e-seguranca-na-aws-com-implementacao-de-landing-zone/logo-autoglass-RGB_LOGO-ORIGINAL-para-fundos-claros-3.png",
        "output_name": "Case Autoglass",
    },
    # ---- CASE 12: Prevent Senior ----
    {
        "headline": "Prevent Senior reduz 18% dos custos mensais com DevOps na AWS.",
        "titulo_secao1": "COMO A PREVENT SENIOR OTIMIZOU SUA INFRAESTRUTURA",
        "texto_cliente": "Especializada em serviços de saúde, a Prevent Senior é dedicada a oferecer soluções inovadoras e de alta qualidade, com abordagem focada na prevenção. Conta com planos de saúde, cuidados médicos, hospitais próprios e centros de diagnóstico, investindo constantemente em modernização e tecnologia avançada.",
        "titulo_secao2": "QUANDO A INOVAÇÃO EM SAÚDE EXIGE EFICIÊNCIA",
        "texto_desafio": "A empresa precisava aprimorar seus sistemas robustos e integrados para gerenciar operações e garantir fluxo eficiente de informações médicas e administrativas. Incorpora tecnologia de ponta incluindo equipamentos avançados de diagnóstico, telemedicina e inteligência artificial.\n\nBuscou a dataRain para implementar DevOps eficiente para otimizar processos de desenvolvimento e implantação, garantindo alta disponibilidade, escalabilidade e segurança.",
        "titulo_secao3": "PARCERIA DATARAIN: DEVOPS PARA SAÚDE",
        "texto_solucao": "A estratégia apostou em três pilares:\nPipeline CI/CD com Jenkins, Bitbucket e Nexus integrados à AWS\nEKS com ArgoCD e Karpenter para containers escaláveis\nRDS Multi-AZ, DynamoDB, Kinesis Streams e ElastiCache para dados",
        "texto_voce_sabia": "A Prevent Senior utiliza integração entre Azure AD e AWS SSO para autenticação unificada. O Vault é utilizado para armazenamento seguro e criptografia de dados sensíveis. Logs são centralizados em conta de auditoria separada para conformidade.",
        "texto_transicao": "Com o desafio de otimizar uma infraestrutura complexa de saúde, a proposta era clara: implementar DevOps para reduzir custos, aumentar agilidade e garantir segurança dos dados médicos.",
        "titulo_arquitetura": "DEVOPS PARA SAÚDE: EFICIÊNCIA E SEGURANÇA",
        "texto_arquitetura": "Pipeline CI com Jenkins, Bitbucket e Nexus. ALB para EC2 e EKS. Docker Registry com ArgoCD para deploy automatizado. Karpenter para dimensionamento eficiente de containers.\n\nRDS Multi-AZ para alta disponibilidade, DynamoDB para NoSQL, Kinesis Streams para dados em tempo real, ElastiCache para cache. Azure AD integrado ao AWS SSO. Vault para criptografia.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: 18% DE ECONOMIA",
        "texto_resultados": "Redução de custos mensais da AWS em 18,18%, com impacto financeiro significativo. Processos de desenvolvimento e implantação otimizados com entrega mais rápida de funcionalidades.\n\nMonitoramento eficaz com Prometheus, Grafana e New Relic. Segmentação de VPC com sub-redes públicas e privadas, NAT Gateways e controle de acesso EKS.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "SAÚDE DIGITAL EFICIENTE",
        "texto_continuacao": "A execução das soluções DevOps resultou em infraestrutura de TI mais eficiente, ágil e segura. Medidas de segurança com IAM, Azure AD e princípio do menor privilégio garantem gerenciamento seguro.\n\nA arquitetura de rede com segmentação de VPC e controles de acesso garante um ambiente seguro e controlado para dados sensíveis de saúde.",
        "logo_path": "public/content-images/cases/implementacao-de-devops-garante-economia-para-a-prevent-senior/2.png.webp",
        "output_name": "Case Prevent Senior",
    },
    # ---- CASE 13: RNP Landing Zone ----
    {
        "headline": "RNP alcança 200% mais eficiência na configuração de contas com Landing Zone na AWS.",
        "titulo_secao1": "COMO A RNP REESTRUTUROU SUA LANDING ZONE",
        "texto_cliente": "A Rede Nacional de Ensino e Pesquisa (RNP) é uma plataforma de comunicação e colaboração digital que promove inovação em tecnologia da informação. Sua atuação abrange o setor público, focando especialmente na área de Educação, conectando instituições de ensino e pesquisa em todo o Brasil.",
        "titulo_secao2": "QUANDO A SEGURANÇA PRECISA EVOLUIR",
        "texto_desafio": "A infraestrutura de nuvem da RNP, estabelecida em 2018, apresentava lacunas significativas em segurança e conformidade. A Landing Zone estava desatualizada, a estrutura de contas era incompatível com o crescimento da organização.\n\nFaltava baseline de segurança, criptografia centralizada de logs do CloudTrail e políticas de controle de serviços (SCPs). A gestão de usuários e acessos era trabalhosa e ineficiente.",
        "titulo_secao3": "PARCERIA DATARAIN: SEGURANÇA REESTRUTURADA",
        "texto_solucao": "A estratégia apostou em três pilares:\nReestruturação da Landing Zone com AWS Control Tower\nImplementação de SCPs e federação Azure ADFS com IAM Identity Center\nConfiguração de serviços nativos de segurança (CloudTrail, Config, GuardDuty)",
        "texto_voce_sabia": "A RNP conecta mais de 800 instituições de ensino e pesquisa no Brasil através do programa NasNuvens. A reestruturação da Landing Zone trouxe 200% mais eficiência na configuração de contas e criação de novas contas 2x mais rápida.",
        "texto_transicao": "Com o desafio de elevar a maturidade de segurança de uma infraestrutura que conecta centenas de instituições, a proposta era clara: reestruturar a Landing Zone seguindo as melhores práticas da AWS.",
        "titulo_arquitetura": "LANDING ZONE: SEGURANÇA E ESCALABILIDADE",
        "texto_arquitetura": "Nova estrutura de contas com OUs para segurança, recursos compartilhados e gestão de usuários. Configuração de Control Tower, CloudTrail, Config, IAM Identity Center e SCPs.\n\nFederação de usuários entre Azure ADFS e AWS IAM Identity Center. Serviços nativos de segurança para garantir alto padrão de conformidade com diretrizes AWS.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: EFICIÊNCIA E GOVERNANÇA",
        "texto_resultados": "200% mais eficiência na configuração de contas com padrão de segurança atualizado. Criação de novas contas 2x mais rápida com configurações pré-definidas via AWS Organizations.\n\nReestruturação da Landing Zone com segurança e escalabilidade. Automatização e simplificação de processos com redução drástica do trabalho operacional.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "MATURIDADE CONTÍNUA",
        "texto_continuacao": "Para elevar ainda mais a maturidade de segurança, a dataRain sugeriu atividades adicionais incluindo políticas de tags e backup, estratégia de Disaster Recovery e outras iniciativas.\n\nA reestruturação garantiu conformidade com exigências regulatórias e políticas de retenção de logs, fortalecendo a segurança de toda a organização.",
        "logo_path": "public/content-images/cases/implementacao-de-landingzone-em-ambiente-de-nuvem-aws-na-rnp/rnp-logo.png",
        "output_name": "Case RNP Landing Zone",
    },
    # ---- CASE 14: Sírio-Libanês ----
    {
        "headline": "Hospital Sírio-Libanês integra dados médicos ao padrão global FHIR com serverless na AWS.",
        "titulo_secao1": "COMO O SÍRIO-LIBANÊS INTEGROU SEUS DADOS MÉDICOS",
        "texto_cliente": "O Hospital Sírio-Libanês, fundado em 1921, é reconhecido por sua excelência médica e compromisso com a qualidade dos cuidados no Brasil. Com equipe altamente qualificada e infraestrutura moderna, oferece ampla gama de especialidades médicas, garantindo diagnósticos precisos e tratamentos eficazes.",
        "titulo_secao2": "QUANDO SISTEMAS DIFERENTES PRECISAM CONVERSAR",
        "texto_desafio": "Na jornada de modernização, a integração entre diferentes sistemas é um desafio crítico. Cada sistema possui seu próprio padrão de tratamento e armazenamento de informações, levando a problemas de comunicação e interoperabilidade.\n\nO hospital enfrentava o desafio de adaptar seu padrão interno (protocolo HSL) ao padrão global FHIR (Fast Healthcare Interoperability Resources), permitindo comunicação eficaz entre os sistemas.",
        "titulo_secao3": "PARCERIA DATARAIN: INTEROPERABILIDADE EM SAÚDE",
        "texto_solucao": "A estratégia apostou em três pilares:\nInterface bidirecional para mapear e converter protocolo HSL para FHIR\nAWS Lambda para processamento serverless dos dados\nAWS HealthLake para armazenar e validar arquivos FHIR",
        "texto_voce_sabia": "FHIR (Fast Healthcare Interoperability Resources) é o padrão global para troca de dados de saúde. A integração permite que hospitais compartilhem informações médicas de forma padronizada, melhorando a continuidade do cuidado ao paciente.",
        "texto_transicao": "Com o desafio de integrar um sistema proprietário ao padrão global de saúde, a proposta era clara: criar uma interface serverless que convertesse dados de forma eficaz e segura.",
        "titulo_arquitetura": "INTEGRAÇÃO FHIR: SERVERLESS E SEGURA",
        "texto_arquitetura": "A interface mapeia componentes do protocolo HSL e converte para o padrão FHIR, estabelecendo canal bidirecional de comunicação. AWS Lambda processa e converte os dados de forma serverless.\n\nAWS HealthLake armazena, consulta e valida os arquivos FHIR gerados, assegurando armazenamento seguro e em conformidade com padrões de saúde. Desenvolvido em JavaScript com Node.js.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: DADOS INTEGRADOS",
        "texto_resultados": "Integração aprimorada entre protocolo interno e padrão global FHIR. Eficiência operacional com processamento serverless via AWS Lambda. Armazenamento seguro e conforme com AWS HealthLake.\n\nAgilidade e flexibilidade com JavaScript e Node.js. Comunicação eficaz e confiável entre diferentes sistemas de saúde.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "INOVAÇÃO EM SAÚDE DIGITAL",
        "texto_continuacao": "O hospital continua comprometido com a evolução de sua infraestrutura e capacidade de integração. Os próximos passos incluem aprimoramento contínuo da interface e exploração de novas tecnologias.\n\nA parceria entre o Sírio-Libanês, a dataRain e a AWS permitiu superar desafios de integração, fortalecendo a capacidade de oferecer cuidados de saúde de qualidade.",
        "logo_path": "public/content-images/cases/integracao-eficiente-de-dados-medicos-no-hospital-sirio-libanes/HSL_PH_rgb_pos.jpg",
        "output_name": "Case Sirio Libanes",
    },
    # ---- CASE 15: Intelbras ----
    {
        "headline": "Intelbras transforma infraestrutura de energia com DevOps e serverless na AWS.",
        "titulo_secao1": "COMO A INTELBRAS MODERNIZOU SUA INFRAESTRUTURA",
        "texto_cliente": "A Intelbras fornece soluções em segurança, controle de acesso, redes, comunicação, além de produtos e serviços relacionados à geração de energia e aproveitamento da energia solar. Com mais de 40 anos de experiência, é referência no mercado pelo foco em tecnologia e inovação.",
        "titulo_secao2": "QUANDO O CRESCIMENTO EXIGE AUTOMAÇÃO",
        "texto_desafio": "A Intelbras precisava estabelecer uma infraestrutura capaz de acompanhar o crescimento e as demandas do seu negócio de energia. Carecia de uma abordagem que garantisse escalabilidade, segurança e confiabilidade.\n\nBuscou implementação de DevOps para estabelecer processos automatizados de desenvolvimento, entrega e operação, agilizando o ciclo de vida de seus produtos e serviços.",
        "titulo_secao3": "PARCERIA DATARAIN: DEVOPS PARA ENERGIA",
        "texto_solucao": "A estratégia apostou em três pilares:\nAbordagem serverless com AWS Lambda, Fargate e Serverless Framework\nTestes automatizados com Vitest em Node.js no pipeline\nMonitoramento personalizado com CloudWatch e Terraform para infraestrutura",
        "texto_voce_sabia": "A Intelbras utiliza Serverless Framework com CloudFormation para infraestrutura como código, e pipeline GitLab para gerenciar alterações, construir infraestrutura, realizar testes e publicar código. A arquitetura abrange três zonas de disponibilidade na região de São Paulo.",
        "texto_transicao": "Com o desafio de modernizar a infraestrutura de energia de uma empresa com mais de 40 anos, a proposta era clara: implementar DevOps com serverless para máxima eficiência e automação.",
        "titulo_arquitetura": "DEVOPS + SERVERLESS: AUTOMAÇÃO COMPLETA",
        "texto_arquitetura": "AWS Lambda e Fargate para processamento serverless. Serverless Framework com CloudFormation para infraestrutura como código. Pipeline GitLab para CI/CD automatizado com testes Vitest.\n\nIdentity Center (SSO) para autenticação. CloudWatch para monitoramento. Terraform para mudanças de infraestrutura. Criptografia em RDS e DynamoDB com KMS. RPO de 24h e RTO definidos.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: EFICIÊNCIA OPERACIONAL",
        "texto_resultados": "Gestão eficiente da infraestrutura com Serverless Framework, Lambda e Fargate. Processo de implantação automatizado com maior agilidade na entrega de novas versões.\n\nAlta disponibilidade e escalabilidade com três zonas de disponibilidade. Gerenciamento de configuração e monitoramento aprimorados com identificação proativa de problemas.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "INOVAÇÃO EM ENERGIA",
        "texto_continuacao": "A parceria estratégica com a dataRain assegurou à Intelbras a capacidade de se adaptar rapidamente às demandas do mercado, fornecendo soluções inovadoras e de alto desempenho.\n\nA automação proporcionada pelo DevOps permite provisionar e gerenciar recursos de maneira simplificada, garantindo disponibilidade e otimizando o uso dos serviços em nuvem.",
        "logo_path": "public/content-images/cases/intelbras-devops-transforma-infraestrutura-e-impulsiona-eficiencia-operacional/intelbras-logo.png",
        "output_name": "Case Intelbras",
    },
    # ---- CASE 16: Seade ----
    {
        "headline": "Fundação Seade reduz custos de US$ 5 milhões para US$ 0,3 milhão migrando para a AWS.",
        "titulo_secao1": "COMO O SEADE MODERNIZOU SUA INFRAESTRUTURA",
        "texto_cliente": "A Fundação Seade é um órgão da Secretaria de Planejamento e Gestão do Governo de São Paulo, centro de referência nacional na produção e disseminação de análises e estatísticas socioeconômicas e demográficas. Ao longo de 40 anos, tem se constituído em fonte segura e atualizada de dados sobre o estado de São Paulo.",
        "titulo_secao2": "QUANDO OS CUSTOS INVIABILIZAM A MODERNIZAÇÃO",
        "texto_desafio": "Em 2018, o volume de dados estava crescendo a velocidade sem precedentes, sendo necessário redesenhar a infraestrutura tecnológica. Os custos de modernização do parque tecnológico eram proibitivos: US$ 5 milhões.\n\nServiços de hospedagem custariam US$ 1,25 milhão. A licitação pública resultou em um plano de migração IaaS na AWS com a dataRain custando US$ 0,3 milhão — 1/60 do CAPEX necessário.",
        "titulo_secao3": "PARCERIA DATARAIN: MIGRAÇÃO ECONÔMICA",
        "texto_solucao": "A estratégia apostou em três pilares:\nMigração IaaS para AWS com custo de US$ 0,3 milhão (1/60 do CAPEX)\nAmazon EC2, EBS e S3 para computação e armazenamento\nConsultoria, gerenciamento de ambiente e operação DevOps",
        "texto_voce_sabia": "A Fundação Seade dissemina sua produção com internet como principal instrumento, gratuidade de acesso e disponibilização de todo o acervo recente. A migração para a nuvem garantiu escalabilidade para acompanhar o crescimento contínuo dos dados socioeconômicos.",
        "texto_transicao": "Com o desafio de modernizar uma infraestrutura de dados com custos proibitivos, a proposta era clara: migrar para a nuvem AWS com economia radical e escalabilidade garantida.",
        "titulo_arquitetura": "MIGRAÇÃO ECONÔMICA: DO LEGADO À NUVEM",
        "texto_arquitetura": "A infraestrutura utiliza Amazon EC2 para computação, Amazon EBS para armazenamento em bloco, Amazon S3 para armazenamento de objetos, Amazon Route 53 para DNS e Amazon VPC para rede isolada.\n\nServiços agregados incluem setup, consultoria, gerenciamento de ambiente e operação DevOps contínua.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: ECONOMIA DE 94%",
        "texto_resultados": "Redução de custos de US$ 5 milhões (CAPEX) para US$ 0,3 milhão na nuvem AWS — economia de 94%. Infraestrutura escalável para acompanhar o crescimento do volume de dados.\n\nA dataRain se tornou trusted advisor do Seade, demonstrando profissionalismo e transparência. A AWS e a dataRain ajudam com consultoria de TI e serviços em nuvem em qualquer desafio.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "DADOS PÚBLICOS NA NUVEM",
        "texto_continuacao": "A migração garantiu que o Seade continue sendo fonte segura e atualizada de dados sobre São Paulo, com infraestrutura moderna e escalável para suportar o crescimento contínuo.\n\nA parceria com a dataRain demonstrou que a modernização do setor público é possível com economia radical e qualidade de serviço.",
        "logo_path": "public/content-images/cases/modernizacao-tecnologica-com-maior-rapidez-no-acesso-aos-dados/8.png.webp",
        "output_name": "Case Seade",
    },
    # ---- CASE 17: Unicamp ----
    {
        "headline": "Unicamp alcança 40% de aumento na eficiência operacional com migração para AWS.",
        "titulo_secao1": "COMO A UNICAMP TRANSFORMOU SUA GESTÃO DE SAÚDE",
        "texto_cliente": "A Universidade Estadual de Campinas (Unicamp), fundada em 1966, é reconhecida nacional e internacionalmente por sua excelência acadêmica e contribuições significativas para o desenvolvimento científico, tecnológico e cultural do país. O sistema AGHUse gerencia suas operações de saúde.",
        "titulo_secao2": "QUANDO A INFRAESTRUTURA LIMITA A SAÚDE",
        "texto_desafio": "A infraestrutura padrão do AGHUse apresentava limitações de escalabilidade e custos operacionais elevados. A instituição lidava com crescentes demandas por serviços de saúde e restrições tecnológicas.\n\nA necessidade era clara: migração completa do serviço com melhoria de performance e redução de custos de operação e armazenamento de 150 TB de dados, incluindo imagens PACS.",
        "titulo_secao3": "PARCERIA DATARAIN: SAÚDE NA NUVEM",
        "texto_solucao": "A estratégia apostou em três pilares:\nMigração sem downtime para VMs escaláveis e otimizadas na AWS\nAmazon S3 para armazenamento seguro e escalável de 150 TB\nSegurança avançada com GuardDuty, Inspector e IAM",
        "texto_voce_sabia": "A migração do sistema AGHUse foi realizada sem nenhuma interrupção nos serviços de saúde da Unicamp. A solução utiliza VPN e firewalls para proteger dados contra acessos não autorizados, com integrações e automações serverless.",
        "texto_transicao": "Com o desafio de migrar 150 TB de dados de saúde sem interrupção dos serviços, a proposta era clara: criar uma infraestrutura segura, escalável e moderna na AWS.",
        "titulo_arquitetura": "SAÚDE NA NUVEM: SEGURA E ESCALÁVEL",
        "texto_arquitetura": "VMs otimizadas para performance e segurança do AGHUse. Amazon S3 para armazenamento seguro e escalável dos dados dos pacientes. Rede privada com VPN e firewalls.\n\nSegurança avançada com Amazon GuardDuty, Inspector e IAM. Integrações e automações serverless para eficiência operacional. Migração sem downtime com técnicas avançadas.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: EFICIÊNCIA E ECONOMIA",
        "texto_resultados": "40% de aumento na eficiência operacional. 30% de redução de custos totais. 25% de salto na disponibilidade geral do sistema AGHUse.\n\nMigração sem interrupções nos serviços de saúde. Infraestrutura segura e escalável para 150 TB de dados médicos.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "FUTURO DA SAÚDE DIGITAL",
        "texto_continuacao": "A Unicamp avalia expandir a nuvem AWS para pesquisa e ensino, potencializando eficiência e inovação em toda a instituição. Explora inteligência artificial para diagnósticos, blockchain para segurança de dados e IoT para monitoramento de pacientes.\n\nA parceria com a dataRain proporcionou suporte técnico e estratégico durante todo o processo de transformação digital.",
        "logo_path": "public/content-images/cases/nuvem-aws-impulsiona-eficiencia-da-unicamp-em-40-e-transforma-gestao-da-saude/Logo_Unicamp__0.jpg",
        "output_name": "Case Unicamp",
    },
    # ---- CASE 18: INC Cardiologia ----
    {
        "headline": "INC analisa dados genéticos de 700+ pacientes com IA na AWS para prever riscos cardiovasculares.",
        "titulo_secao1": "COMO O INC REVOLUCIONOU A CARDIOLOGIA COM IA",
        "texto_cliente": "O Instituto Nacional de Cardiologia (INC), localizado no Rio de Janeiro, é uma das instituições mais respeitadas no campo da saúde cardiovascular no Brasil. Como parte do projeto IPOLE, o INC colabora com especialistas internacionais para enfrentar desafios críticos da área da saúde.",
        "titulo_secao2": "QUANDO DADOS MULTIMODAIS EXIGEM IA AVANÇADA",
        "texto_desafio": "O INC enfrentava dificuldades na análise de dados multimodais (genéticos, clínicos e de imagem) de mais de 700 pacientes brasileiros do estudo Renomica. Faltava infraestrutura escalável para armazenar e processar grandes quantidades de dados.\n\nA correlação entre dados genéticos e fenotípicos exigia modelos avançados de IA e ML. A manipulação de dados sensíveis deveria estar em conformidade com a LGPD.",
        "titulo_secao3": "PARCERIA DATARAIN: IA PARA CARDIOLOGIA",
        "texto_solucao": "A estratégia apostou em três pilares:\nAWS SageMaker para análise exploratória e treinamento de modelos de ML\nAWS HealthOmics para gestão e processamento de dados genéticos\nAmazon Bedrock para NLP e análise avançada de imagens médicas",
        "texto_voce_sabia": "O projeto utilizou SageMaker + Bedrock Multimodal Embeddings para visualização combinada de dados genéticos e de imagem, permitindo identificar padrões correlacionados entre genoma e fenótipo para previsão de riscos cardiovasculares.",
        "texto_transicao": "Com o desafio de analisar dados genéticos e fenotípicos de centenas de pacientes, a proposta era clara: usar IA avançada na AWS para identificar biomarcadores digitais e prever riscos cardiovasculares.",
        "titulo_arquitetura": "IA EM CARDIOLOGIA: GENÔMICA + IMAGEM",
        "texto_arquitetura": "Fase 1: AWS Storage Gateway para transferência segura, S3 para armazenamento, SageMaker para análise e treinamento de modelos, HealthOmics para dados genéticos.\n\nFase 2: S3 para imagens DICOM (ressonâncias e ECGs), Bedrock para NLP e análise de imagens, SageMaker + Bedrock Multimodal Embeddings para visualização combinada de dados genéticos e de imagem.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: BIOMARCADORES DIGITAIS",
        "texto_resultados": "Melhor organização e acesso aos dados genéticos e fenotípicos de 700+ pacientes. Modelos preditivos para estratificação de pacientes com base em características genéticas e clínicas.\n\nAmazon Bedrock melhorou interpretação de relatórios médicos e extração de padrões clínicos. Conformidade com LGPD via anonimização e controle de acessos.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "FUTURO DA CARDIOLOGIA PREDITIVA",
        "texto_continuacao": "Os próximos passos incluem expandir o projeto para novos pacientes, criar pipelines automatizados para coleta e análise de dados, e publicar resultados em artigos científicos.\n\nA adoção de AWS HealthLake para armazenamento integrado de dados de saúde está sendo explorada para ampliar as capacidades da solução.",
        "logo_path": "public/content-images/cases/revolucao-na-cardiologia-analise-de-dados-geneticos-e-fenotipicos-com-ia-na-aws/logoINC.jpg",
        "output_name": "Case INC Cardiologia",
    },
    # ---- CASE 19: RNP Conferência ----
    {
        "headline": "RNP reduz 15% dos custos anuais migrando Conferência Web da Azure para AWS.",
        "titulo_secao1": "COMO A RNP OTIMIZOU SEU SERVIÇO DE CONFERÊNCIA",
        "texto_cliente": "O Conferência Web é um serviço de educação à distância da RNP, desenvolvido para oferecer experiência completa e segura em interações online e trabalho colaborativo. Combina vídeo e áudio para criação de salas virtuais voltadas para aulas, reuniões, palestras e projetos.",
        "titulo_secao2": "QUANDO A NUVEM ANTERIOR NÃO ATENDE",
        "texto_desafio": "O desafio era reduzir custos com a nuvem. Inicialmente com ambiente on-premises, a RNP migrou para a Azure, mas buscava mais qualidade por menos custo. Após prova de conceito, a AWS foi validada principalmente pelo quesito economia.\n\nA plataforma precisava sustentar arquitetura de containers, cluster e Kubernetes, com alta disponibilidade e capacidade ilimitada de armazenamento.",
        "titulo_secao3": "PARCERIA DATARAIN: MIGRAÇÃO AZURE → AWS",
        "texto_solucao": "A estratégia apostou em três pilares:\nMigração da Azure para AWS via programa NasNuvens\nArquitetura de containers e Kubernetes com Amazon EKS\nSolução mais enxuta e barata para a carga de trabalho",
        "texto_voce_sabia": "O programa NasNuvens da RNP conta com provedores globais qualificados para oferta de IaaS às 800 instituições do Sistema RNP. A AWS foi a solução que atendeu todos os requisitos e apresentou o melhor custo-benefício.",
        "texto_transicao": "Com o desafio de reduzir custos e ganhar escalabilidade para um serviço que atende centenas de instituições, a proposta era clara: migrar da Azure para a AWS com arquitetura de containers.",
        "titulo_arquitetura": "CONTAINERS E KUBERNETES NA AWS",
        "texto_arquitetura": "A plataforma foi remodelada para sustentar arquitetura de containers e cluster na AWS, com Amazon EKS para cargas Kubernetes. Integração com pipeline DevOps para entrega contínua.\n\nUtiliza Amazon S3, Lambda, EC2, Route 53 e CloudTrail. Alta disponibilidade e capacidade ilimitada de armazenamento com dados digitalizados, organizados e padronizados.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: 15% DE ECONOMIA ANUAL",
        "texto_resultados": "Economia de 15% ao ano com a nuvem AWS, comprovada por análise TCO. Serviço altamente escalável com velocidade ilimitada e sem interrupção. Melhor integração com ferramentas de pipeline DevOps.\n\nMaturidade, confiabilidade e plataforma aderente à arquitetura de containers, cluster e Kubernetes, todas atendidas pela AWS.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "EDUCAÇÃO DIGITAL EM ESCALA",
        "texto_continuacao": "Com o NasNuvens, a jornada para a nuvem foi mais curta e rápida. Os provedores já estão qualificados no catálogo de serviços, trazendo agilidade, confiabilidade e qualidade na contratação.\n\nA RNP está pronta para oferecer consultoria sobre a melhor arquitetura para cada demanda, já pensando na nuvem e pronta para expandir seus registros.",
        "logo_path": "public/content-images/cases/rnp-reduz-significativamente-custos-migrando-seu-servico-de-conferencia-web-para-da-nuvem-aws/rnp-logo.png",
        "output_name": "Case RNP Conferencia",
    },
    # ---- CASE 20: Desenvolve SP ----
    {
        "headline": "Desenvolve SP migra de colocation para AWS em 45 dias com 20% de economia anual.",
        "titulo_secao1": "COMO O DESENVOLVE SP MIGROU PARA A NUVEM",
        "texto_cliente": "O Desenvolve SP é o Banco do Empreendedor, instituição financeira do Estado de São Paulo que oferece financiamento para pequenas e médias empresas. Em 2021, desembolsou R$ 736,1 milhões, dos quais R$ 186,23 milhões para créditos sustentáveis, um aumento de 81,7% sobre o ano anterior.",
        "titulo_secao2": "QUANDO O COLOCATION LIMITA O CRESCIMENTO",
        "texto_desafio": "O banco usava colocation de outra provedora, que não trazia ambiente seguro exigido pela instituição bancária e tinha custos muito altos, sem clareza na apuração dos valores. Qualquer crescimento demandava equipamentos físicos e processo licitatório caro e lento.\n\nEm busca de mais segurança e economia, o Desenvolve SP decidiu migrar para a nuvem AWS para ter mais controle sobre informações e custos, além de escalar conforme a demanda.",
        "titulo_secao3": "PARCERIA DATARAIN: MIGRAÇÃO EM 45 DIAS",
        "texto_solucao": "A estratégia apostou em três pilares:\nMigração completa em apenas 45 dias com preço fixo\nRedução de risco com processos definidos e repetíveis\nAmpla gama de serviços AWS para segurança e conformidade bancária",
        "texto_voce_sabia": "O Desenvolve SP utiliza 17 serviços AWS incluindo S3, EC2, Security Hub, Systems Manager, Secrets Manager, Glue, Config, KMS, CloudTrail e Cost Explorer. A migração de preço fixo eliminou gastos duplos e reduziu prazos.",
        "texto_transicao": "Com o desafio de sair de um colocation caro e inseguro para a nuvem, a proposta era clara: migrar rapidamente com economia, segurança e transparência de custos.",
        "titulo_arquitetura": "DO COLOCATION À NUVEM AWS",
        "texto_arquitetura": "A infraestrutura utiliza Amazon S3, EC2, Route 53, VPC, CloudWatch e CloudTrail como base. Segurança com Security Hub, Config, KMS, Secrets Manager e Certificate Manager.\n\nSystems Manager para gerenciamento, Glue para ETL, SNS e SQS para mensageria, e Cost Explorer para visibilidade de custos. Migração completa em 45 dias.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: ECONOMIA E TRANSPARÊNCIA",
        "texto_resultados": "20% de economia no custo anual com Data Center. Previsibilidade de custos e clareza de informações sobre o serviço contratado. Melhor gestão do crescimento e planejamento.\n\nMigração rápida e segura em 45 dias, levando todos os sistemas. Modernização do ambiente com ajuste rápido e transparente.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "BANCO DIGITAL DO EMPREENDEDOR",
        "texto_continuacao": "A migração para a AWS trouxe ao Desenvolve SP a segurança exigida por uma instituição bancária, com custos transparentes e previsíveis. A escalabilidade da nuvem acompanha o crescimento dos desembolsos de crédito.\n\nA confiança na dataRain, já prestadora de serviços da Prodesp, foi fundamental para conduzir o processo de migração com sucesso.",
        "logo_path": "public/content-images/cases/sair-do-modelo-de-colocation-sem-possibilidade-de-escalar-e-alto-custo-para-nuvem-aws-de-forma-otimizada/11.png.webp",
        "output_name": "Case Desenvolve SP",
    },
    # ---- CASE 21: Setor Elétrico (ANEEL) ----
    {
        "headline": "ANEEL e Instituto ABRADEE criam Sistema de Inteligência Analítica na AWS.",
        "titulo_secao1": "COMO O SETOR ELÉTRICO GANHOU INTELIGÊNCIA ANALÍTICA",
        "texto_cliente": "A ANEEL é a autarquia federal responsável por regular e fiscalizar o setor elétrico no Brasil. O Instituto ABRADEE de Energia atua na promoção de estudos e projetos voltados ao desenvolvimento do setor. Juntos, criaram o Sistema de Inteligência Analítica do Setor Elétrico.",
        "titulo_secao2": "QUANDO DADOS REGULATÓRIOS EXIGEM AGILIDADE",
        "texto_desafio": "O projeto exigia mobilização rápida de toda a infraestrutura para ambientes de Homologação e Produção. Por envolver órgãos reguladores e dados sensíveis do setor elétrico, a solução precisava garantir agilidade, segurança e isolamento.\n\nEscalabilidade para suportar o crescimento do volume de dados analíticos e confiabilidade com alta disponibilidade eram requisitos fundamentais.",
        "titulo_secao3": "PARCERIA DATARAIN: INFRAESTRUTURA REGULATÓRIA",
        "texto_solucao": "A estratégia apostou em três pilares:\nProvisionamento rápido de ambientes de Homologação e Produção\nAmazon EC2, EBS e S3 para processamento e armazenamento analítico\nVPC para isolamento de rede e controle de acesso",
        "texto_voce_sabia": "O Sistema de Inteligência Analítica consolida, processa e disponibiliza dados estratégicos sobre o mercado de energia no Brasil. A cooperação inédita entre ANEEL e Instituto ABRADEE representa um marco na transparência do setor elétrico.",
        "texto_transicao": "Com o desafio de criar uma plataforma analítica para o setor elétrico brasileiro, a proposta era clara: provisionar infraestrutura segura e escalável na AWS dentro do cronograma institucional.",
        "titulo_arquitetura": "INTELIGÊNCIA ANALÍTICA NA NUVEM",
        "texto_arquitetura": "Amazon EC2 para processamento analítico, EBS para armazenamento de alta performance, S3 para datasets e backups, Route 53 para DNS e VPC para isolamento de rede.\n\nServiços agregados incluem setup, consultoria em arquitetura de nuvem, gerenciamento de ambiente e operação DevOps contínua.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: ENTREGA NO PRAZO",
        "texto_resultados": "Ambientes de Homologação e Produção entregues dentro do prazo estipulado. Infraestrutura escalável preparada para o crescimento do volume de dados do setor elétrico.\n\nOperação estável e segura atendendo aos requisitos regulatórios da ANEEL. Redução da complexidade operacional com serviços gerenciados da AWS.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "ENERGIA E DADOS",
        "texto_continuacao": "O Sistema de Inteligência Analítica representa um avanço na transparência e eficiência do setor elétrico brasileiro. A plataforma permite análises estratégicas sobre o mercado de energia.\n\nA parceria com a dataRain garantiu entrega dentro do prazo com infraestrutura segura e escalável, atendendo aos requisitos regulatórios de ambas as instituições.",
        "logo_path": "public/content-images/cases/sistema-de-inteligencia-analitica-do-setor-eletrico/1.png.webp",
        "output_name": "Case Setor Eletrico",
    },
    # ---- CASE 22: InRad ----
    {
        "headline": "InRad automatiza pré-laudos radiológicos com IA Generativa na AWS.",
        "titulo_secao1": "COMO O INRAD AUTOMATIZOU SEUS PRÉ-LAUDOS",
        "texto_cliente": "O Instituto de Radiologia do Hospital das Clínicas da FMUSP (InRad) é centro de referência em diagnóstico por imagem e oncologia terapêutica. Realiza cerca de 330 mil exames de alta complexidade anualmente, com estruturas como o CEDIM e o NDI para digitalização e distribuição de imagens.",
        "titulo_secao2": "QUANDO O VOLUME DE EXAMES EXIGE AUTOMAÇÃO",
        "texto_desafio": "O processo de elaboração de laudos envolvia muitas etapas manuais, dificultando escalabilidade e adoção. A primeira iteração contou com apenas 1000 pacientes, insuficiente para validação robusta do modelo de IA.\n\nA integração com o sistema MV (prontuário eletrônico) estava em definição, e os pré-laudos gerados não estavam suficientemente detalhados, demandando aprimoramentos.",
        "titulo_secao3": "PARCERIA DATARAIN: IA GENERATIVA PARA RADIOLOGIA",
        "texto_solucao": "A estratégia apostou em três pilares:\nAmazon Bedrock com modelo LLM para geração de sumários clínicos\nAWS Step Functions e Lambda para processamento e anonimização\nIntegração com PEP e PACS para fluxo clínico completo",
        "texto_voce_sabia": "83,3% dos sumários gerados receberam nota 4 ou 5 (máxima) e 93,3% foram considerados livres de alucinações. A solução economiza tempo dos radiologistas, que visualizam informações organizadas de forma mais rápida e estruturada.",
        "texto_transicao": "Com o desafio de automatizar a geração de pré-laudos para 330 mil exames anuais, a proposta era clara: usar IA Generativa na AWS para elevar produtividade e padronização.",
        "titulo_arquitetura": "IA GENERATIVA: DO PRONTUÁRIO AO PRÉ-LAUDO",
        "texto_arquitetura": "Integração com PEP e PACS para ingestão de documentos do MV. Amazon S3 para armazenamento, Step Functions e Lambda para processamento e anonimização por paciente.\n\nAmazon Bedrock com modelo LLM gera sumário clínico. Médicos acessam pré-laudos via link no S3. Dados removidos do ambiente AWS após validação para conformidade com segurança hospitalar.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: 83,3% DE NOTAS MÁXIMAS",
        "texto_resultados": "83,3% dos sumários gerados receberam nota 4 ou 5 (máxima). 93,3% dos sumários foram considerados livres de alucinações, erro comum em modelos de IA.\n\nEconomia de tempo para os radiologistas, que passaram a visualizar informações organizadas de forma mais rápida e estruturada.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "RADIOLOGIA DO FUTURO",
        "texto_continuacao": "A dataRain liderou a implementação e evolução da solução, garantindo maior robustez técnica e adequação ao fluxo clínico do hospital. A IA Generativa agrega valor ao trabalho dos radiologistas.\n\nA solução utiliza diferentes fontes de dados — formulários médicos, pedidos de exame e históricos clínicos — para otimizar e padronizar a elaboração de laudos radiológicos.",
        "logo_path": "public/content-images/cases/transformacao-digital-no-inrad-automacao-de-pre-laudos-radiologicos-com-ia-generativa-e-aws/logo-inrad.png",
        "output_name": "Case InRad",
    },
    # ---- CASE 23: IPDA ----
    {
        "headline": "IPDA alcança 99,99% de disponibilidade e 500% mais velocidade com AWS.",
        "titulo_secao1": "COMO A IPDA MANTEVE SEU SITE NA PANDEMIA",
        "texto_cliente": "A Igreja Pentecostal Deus É Amor (IPDA), fundada em 1962, conta com mais de 17 mil templos no Brasil e presença em 88 países, com mais de um milhão de fiéis pelo mundo. Seu principal templo é considerado um dos maiores templos evangélicos do mundo.",
        "titulo_secao2": "QUANDO A PANDEMIA MULTIPLICA O TRÁFEGO",
        "texto_desafio": "Com a quarentena do COVID-19, houve grande aumento de visitantes ao website para acompanhar cultos online. Os acessos saltaram de 150 mil para mais de 310 mil mensais. Como organização sem fins lucrativos, a IPDA possuía recursos financeiros limitados.\n\nPrecisava de solução escalável, rápida de implementar e com elasticidade para ajustar recursos conforme demanda. O evento virtual de aniversário previa tráfego massivo ainda maior.",
        "titulo_secao3": "PARCERIA DATARAIN: PERFORMANCE SOB PRESSÃO",
        "texto_solucao": "A estratégia apostou em três pilares:\nAmazon CloudFront para CDN com aumento de 500% na velocidade\nAuto Scaling para elasticidade automática conforme demanda\nElastiCache para aliviar carga nos servidores de banco de dados",
        "texto_voce_sabia": "Durante o evento virtual de aniversário da igreja, com tráfego massivo previsto, não só houve disponibilidade total dos serviços, como o desempenho do site foi registrado como melhor do que nunca, graças à nova infraestrutura AWS.",
        "texto_transicao": "Com o desafio de manter um website disponível durante picos de tráfego na pandemia, a proposta era clara: implementar infraestrutura elástica e performática na AWS com recursos limitados.",
        "titulo_arquitetura": "ALTA PERFORMANCE: ELÁSTICA E SEGURA",
        "texto_arquitetura": "Amazon VPC para rede, EC2 para computação, Route 53 para DNS, CloudFront para CDN, Auto Scaling para elasticidade, ElastiCache para cache de banco de dados.\n\nIntegração com Active Directory para autenticação. CloudWatch para monitoramento. KMS para criptografia. Lambda para automações. Certificado SSL e regras de firewall configurados.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: DISPONIBILIDADE TOTAL",
        "texto_resultados": "99,99% de disponibilidade do serviço. Aumento de 500% na velocidade de carregamento com CloudFront. Bloqueio de 2x mais tráfego malicioso que o ambiente anterior.\n\nEscalabilidade automática real com tempo de resposta mínimo. Modelo de pagamento sob demanda ideal para organização sem fins lucrativos.",
        "texto_depoimento": '"Trabalhamos muito próximos do time da IPDA, conseguimos sentir a dor deles, calçar os sapatos do cliente. Essa era a única forma de oferecer a experiência da nuvem que eles merecem."',
        "texto_autor_depoimento": "- Diretor de Tecnologia da dataRain",
        "titulo_continuacao": "FÉ E TECNOLOGIA",
        "texto_continuacao": "O modelo de pagamento sob demanda foi um ótimo benefício para a organização sem fins lucrativos, democratizando o acesso às mais novas tecnologias de forma economicamente viável.\n\nA tomada de decisões perante períodos sazonais está mais fácil, com otimização na disponibilidade de serviços e custos. Os próximos passos incluem incorporação do sistema de doações pelo website.",
        "logo_path": "public/content-images/cases/website-em-ambiente-de-alta-performance-disponibilidade-e-velocidade/ipda-logo.png",
        "output_name": "Case IPDA",
    },
    # ---- CASE 24: ABIT ----
    {
        "headline": "ABIT migra website institucional para AWS com transparência e agilidade.",
        "titulo_secao1": "COMO A ABIT MODERNIZOU SUA INFRAESTRUTURA",
        "texto_cliente": "A Associação Brasileira da Indústria Têxtil e de Confecção (ABIT), fundada em 1957, representa 27,5 mil empresas que empregam mais de 1,5 milhão de trabalhadores e geram faturamento anual de US$ 51,58 bilhões. O Brasil é a quinta maior indústria têxtil do mundo.",
        "titulo_secao2": "QUANDO A TRANSPARÊNCIA É ESSENCIAL",
        "texto_desafio": "A organização sem fins lucrativos precisava de transparência sobre gastos com TI não oferecidos pelo provedor anterior. Se limitava a KPIs do Google Analytics sem métricas precisas de infraestrutura.\n\nA ABIT entrega diversos eventos e precisa de agilidade criando hotsites. O provedor anterior levava mais de uma semana para provisionar um usuário FTP. Com limitações orçamentárias, cada incremento precisava caber no orçamento.",
        "titulo_secao3": "PARCERIA DATARAIN: NUVEM TRANSPARENTE",
        "texto_solucao": "A estratégia apostou em três pilares:\nMigração cloud-to-cloud com transparência total de custos\nAgilidade na criação de hotsites e gestão de ativos de marketing\nInfraestrutura moderna com políticas de segurança adequadas",
        "texto_voce_sabia": "A ABIT representa a quinta maior indústria têxtil do mundo, com produção de cerca de 5,1 bilhões de peças de vestuário. A infraestrutura moderna na AWS reflete a mentalidade sólida da instituição, atraindo mais associados.",
        "texto_transicao": "Com o desafio de migrar de um provedor sem transparência para a nuvem AWS, a proposta era clara: oferecer visibilidade de custos, agilidade operacional e segurança para uma organização sem fins lucrativos.",
        "titulo_arquitetura": "NUVEM TRANSPARENTE E ÁGIL",
        "texto_arquitetura": "A solução ABIT Cloud IaaS na AWS foi desenvolvida para oferecer transparência total de custos, métricas precisas de infraestrutura e agilidade na criação de hotsites e gestão de ativos.\n\nPolíticas de segurança adequadas e KPIs transparentes refletem a mentalidade sólida da instituição. A dataRain atua como parceiro de confiança para assuntos de TI e inovação.",
        "titulo_resultados": "RESULTADOS NA PRÁTICA: TRANSPARÊNCIA E AGILIDADE",
        "texto_resultados": "Transparência total sobre gastos com TI e métricas precisas de infraestrutura. Agilidade na criação de hotsites e gestão de ativos de marketing para eventos nacionais e internacionais.\n\nInfraestrutura moderna com políticas de segurança adequadas. A ABIT pode se concentrar em outros aspectos do negócio sem preocupação com TI.",
        "texto_depoimento": None,
        "texto_autor_depoimento": None,
        "titulo_continuacao": "INDÚSTRIA TÊXTIL NA NUVEM",
        "texto_continuacao": "A dataRain deu à ABIT a capacidade de criar a solução certa com recursos, capacidades e suporte para atingir objetivos de curto prazo e criar base para estratégia de nuvem de longo prazo.\n\nO envolvimento permitiu que a ABIT se concentrasse em outros aspectos do negócio, com a dataRain sempre disponível para responder sobre tecnologia e inovação.",
        "logo_path": "public/content-images/cases/website-institucional-instalado-em-ambiente-escalavel-seguro-e-de-alta-disponibilidade/6.png.webp",
        "output_name": "Case ABIT",
    },
]


def generate_case(case_data, template_path="cases_ppt/Case PQ.pptx"):
    """Gera um PPT de case a partir do template."""
    output_path = f"{case_data['output_name']}.pptx"
    prs = Presentation(template_path)

    # ============ SLIDE 1 ============
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        text = shape.text.strip()

        if text.startswith("PQ moderniza"):
            replace_shape_text(shape, case_data["headline"])
        elif text == "COMO A PQ MODERNIZOU A COLETA DE DADOS INDUSTRIAIS":
            replace_shape_text(shape, case_data["titulo_secao1"])
        elif text.startswith("No setor químico"):
            replace_shape_text(shape, case_data["texto_cliente"])
        elif text == "QUANDO A EFICIÊNCIA ESBARRA NOS SISTEMAS LEGADOS":
            replace_shape_text(shape, case_data["titulo_secao2"])
        elif text.startswith("A PQ gerenciava"):
            replace_shape_text(shape, case_data["texto_desafio"])
        elif text.startswith("PARCERIA DATARAIN: REESTRUTURANDO"):
            replace_shape_text(shape, case_data["titulo_secao3"])
        elif text.startswith("A estratégia apostou"):
            replace_shape_text(shape, case_data["texto_solucao"])
        elif text.startswith("O monitoramento em tempo real"):
            replace_shape_text(shape, case_data["texto_voce_sabia"])
        elif text.startswith("Com o desafio de reconstruir"):
            replace_shape_text(shape, case_data["texto_transicao"])

    # Substituir logo do cliente (shape 13)
    logo_path = case_data.get("logo_path", "")
    if logo_path and os.path.exists(logo_path):
        replace_image(slide1, 13, logo_path)

    # ============ SLIDE 2 ============
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        text = shape.text.strip()

        if text.startswith("ARQUITETURA FLUIDA"):
            replace_shape_text(shape, case_data["titulo_arquitetura"])
        elif text.startswith("Tudo  começa") or text.startswith("Tudo começa"):
            replace_shape_text(shape, case_data["texto_arquitetura"])
        elif text.startswith("RESULTADOS NA PRÁTICA"):
            replace_multiformat_text(shape, [
                (case_data["titulo_resultados"], True),
                (case_data["texto_resultados"], False),
            ])

    # ============ SLIDE 3 ============
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        text = shape.text.strip()

        if text.startswith('"A DataRain tem sido') or text.startswith('"A dataRain tem sido'):
            if case_data.get("texto_depoimento"):
                replace_multiformat_text(shape, [
                    (case_data["texto_depoimento"], False),
                    ("\n" + case_data["texto_autor_depoimento"], True),
                ])
            else:
                # Sem depoimento: usar texto genérico sobre a parceria
                replace_multiformat_text(shape, [
                    ('"A dataRain se destacou pela proximidade e colaboração com nosso time, '
                     'garantindo transferência de conhecimento e autonomia para operar '
                     'a nova infraestrutura com confiança."', False),
                    ("\n- " + case_data["output_name"].replace("Case ", ""), True),
                ])
        elif text.startswith("AUTONOMIA E ESCALABILIDADE"):
            replace_multiformat_text(shape, [
                (case_data["titulo_continuacao"], True),
                (case_data["texto_continuacao"], False),
            ])

    prs.save(output_path)
    return output_path


def main():
    template = "cases_ppt/Case PQ.pptx"
    if not os.path.exists(template):
        print(f"❌ Template não encontrado: {template}")
        sys.exit(1)

    os.makedirs("cases_ppt", exist_ok=True)

    total = len(CASES)
    sucesso = 0
    erros = []

    for i, case_data in enumerate(CASES, 1):
        nome = case_data["output_name"]
        try:
            output = generate_case(case_data, template)
            # Mover para pasta cases_ppt
            dest = os.path.join("cases_ppt", os.path.basename(output))
            os.rename(output, dest)
            print(f"  [{i:2d}/{total}] ✅ {nome}")
            sucesso += 1
        except Exception as e:
            print(f"  [{i:2d}/{total}] ❌ {nome}: {e}")
            erros.append((nome, str(e)))

    print(f"\n{'='*50}")
    print(f"Total: {sucesso}/{total} gerados com sucesso")
    if erros:
        print(f"Erros ({len(erros)}):")
        for nome, err in erros:
            print(f"  - {nome}: {err}")
    print(f"Arquivos em: cases_ppt/")


if __name__ == "__main__":
    main()
