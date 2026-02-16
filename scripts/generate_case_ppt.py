#!/usr/bin/env python3
"""
Gera um PPT de case a partir do template 'Case PQ.pptx'.
Substitui os textos dos shapes mantendo a formatação original.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import copy
import sys
import os


def replace_shape_text(shape, new_text):
    """Substitui o texto de um shape mantendo a formatação do primeiro run."""
    if not shape.has_text_frame:
        return
    
    tf = shape.text_frame
    
    # Se o novo texto tem múltiplos parágrafos (separados por \n)
    new_paragraphs = new_text.split('\n')
    
    # Pega a formatação do primeiro run do primeiro parágrafo
    first_para = tf.paragraphs[0]
    if not first_para.runs:
        return
    
    ref_run = first_para.runs[0]
    ref_font_bold = ref_run.font.bold
    ref_font_size = ref_run.font.size
    ref_font_color = ref_run.font.color.rgb if ref_run.font.color and ref_run.font.color.rgb else None
    ref_font_name = ref_run.font.name
    ref_para_alignment = first_para.alignment
    ref_para_space_before = first_para.space_before
    ref_para_space_after = first_para.space_after
    ref_para_line_spacing = first_para.line_spacing
    
    # Limpa todos os parágrafos existentes
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    
    # Limpa os runs do primeiro parágrafo
    first_p = tf.paragraphs[0]
    for r in first_p._p.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}r'):
        first_p._p.remove(r)
    
    # Agora adiciona o novo texto
    for idx, para_text in enumerate(new_paragraphs):
        if idx == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        
        para.alignment = ref_para_alignment
        if ref_para_space_before is not None:
            para.space_before = ref_para_space_before
        if ref_para_space_after is not None:
            para.space_after = ref_para_space_after
        if ref_para_line_spacing is not None:
            para.line_spacing = ref_para_line_spacing
        
        run = para.add_run()
        run.text = para_text
        run.font.bold = ref_font_bold
        run.font.size = ref_font_size
        if ref_font_color:
            run.font.color.rgb = ref_font_color
        if ref_font_name:
            run.font.name = ref_font_name


def replace_multiformat_text(shape, sections):
    """
    Substitui texto em shape que tem múltiplas seções com formatação diferente.
    sections = [(text, is_title), ...] onde is_title indica se usa formatação de título.
    Pega a formatação real dos parágrafos existentes.
    """
    if not shape.has_text_frame:
        return
    
    tf = shape.text_frame
    
    # Captura formatação do título (primeiro parágrafo) e corpo (segundo parágrafo)
    title_fmt = {}
    body_fmt = {}
    
    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.bold and run.font.size and run.font.size >= 170000:
                if not title_fmt:
                    title_fmt = {
                        'bold': run.font.bold,
                        'size': run.font.size,
                        'color': run.font.color.rgb if run.font.color and run.font.color.rgb else None,
                        'name': run.font.name,
                        'alignment': para.alignment,
                    }
            else:
                if not body_fmt:
                    body_fmt = {
                        'bold': run.font.bold,
                        'size': run.font.size,
                        'color': run.font.color.rgb if run.font.color and run.font.color.rgb else None,
                        'name': run.font.name,
                        'alignment': para.alignment,
                    }
    
    if not title_fmt:
        title_fmt = body_fmt
    if not body_fmt:
        body_fmt = title_fmt
    
    # Limpa parágrafos
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    
    first_p = tf.paragraphs[0]
    for r in first_p._p.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}r'):
        first_p._p.remove(r)
    
    first = True
    for text, is_title in sections:
        lines = text.split('\n')
        for line in lines:
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            
            fmt = title_fmt if is_title else body_fmt
            para.alignment = fmt.get('alignment')
            
            run = para.add_run()
            run.text = line
            run.font.bold = fmt.get('bold')
            run.font.size = fmt.get('size')
            if fmt.get('color'):
                run.font.color.rgb = fmt.get('color')
            if fmt.get('name'):
                run.font.name = fmt.get('name')


def replace_image(slide, shape_id, new_image_path):
    """Substitui a imagem de um shape mantendo posição e tamanho."""
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    import os
    
    for shape in slide.shapes:
        if shape.shape_id == shape_id and shape.shape_type == 13:
            # Salva posição e tamanho
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            
            # Remove o shape antigo
            sp = shape._element
            sp.getparent().remove(sp)
            
            # Adiciona nova imagem
            pic = slide.shapes.add_picture(new_image_path, left, top, width, height)
            return True
    return False


def generate_autoglass_case():
    template_path = "cases_ppt/Case PQ.pptx"
    output_path = "cases_ppt/Case Autoglass.pptx"
    
    prs = Presentation(template_path)
    
    # ============================================================
    # SLIDE 1
    # ============================================================
    slide1 = prs.slides[0]
    
    for shape in slide1.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        
        text = shape.text.strip()
        
        # Shape 25 - Headline da capa (texto branco grande)
        if text.startswith("PQ moderniza"):
            replace_shape_text(shape,
                "Grupo Autoglass eleva governança e segurança na AWS com implementação de Landing Zone.")
        
        # Shape 22 - Título seção 1
        elif text == "COMO A PQ MODERNIZOU A COLETA DE DADOS INDUSTRIAIS":
            replace_shape_text(shape,
                "COMO O GRUPO AUTOGLASS ELEVOU SUA GOVERNANÇA NA AWS")
        
        # Shape 21 - Texto sobre o cliente
        elif text.startswith("No setor químico"):
            replace_shape_text(shape,
                "O Grupo Autoglass é um ecossistema de soluções que vai além do setor automotivo. "
                "Com a Autoglass, referência em vidros e peças automotivas presente em mais de 85 lojas "
                "no Brasil e na Colômbia, o grupo expandiu sua atuação com a Maxpar, especializada em "
                "assistências para seguradoras, e a B4B.tech, unidade de tecnologia focada em digitalização "
                "e BPO de operações de seguros. A diversidade e a autonomia de cada unidade de negócio "
                "motivaram a busca por uma infraestrutura mais moderna, segura e escalável.")
        
        # Shape 23 - Título seção 2 (Desafio)
        elif text == "QUANDO A EFICIÊNCIA ESBARRA NOS SISTEMAS LEGADOS":
            replace_shape_text(shape,
                "QUANDO O CRESCIMENTO ESBARRA NA GOVERNANÇA")
        
        # Shape 19 - Texto do desafio
        elif text.startswith("A PQ gerenciava"):
            replace_shape_text(shape,
                "O Grupo operava com cerca de seis contas AWS compartilhadas entre cinco Business Units, "
                "o que gerava dificuldades para garantir isolamento entre recursos e segurança das aplicações. "
                "A falta de controle sobre os gastos de cada BU dificultava a alocação de custos e a governança "
                "financeira. A infraestrutura apresentava controles de segurança insuficientes: ausência de "
                "controle sobre regiões AWS utilizadas, descentralização dos logs de segurança e falta de "
                "políticas preventivas e detectivas no nível da organização.\n\n"
                "O Grupo buscava uma solução que proporcionasse autonomia e flexibilidade para cada BU, "
                "sem comprometer a segurança e a rastreabilidade de custos.")
        
        # Shape 24 - Título seção 3 (Solução)
        elif text.startswith("PARCERIA DATARAIN: REESTRUTURANDO"):
            replace_shape_text(shape,
                "PARCERIA DATARAIN: IMPLEMENTANDO A LANDING ZONE")
        
        # Shape 17 - Texto da solução/estratégia
        elif text.startswith("A estratégia apostou"):
            replace_shape_text(shape,
                "A estratégia apostou em três pilares:\n"
                "Segmentar o ambiente com OUs dedicadas, garantindo isolamento e autonomia por BU\n"
                "Configurar serviços nativos de segurança como CloudTrail, GuardDuty e Security Hub\n"
                "Capacitar o time interno com transferência de conhecimento e trabalho colaborativo")
        
        # Shape 20 - Texto de transição/resumo
        elif text.startswith("Com o desafio de reconstruir"):
            replace_shape_text(shape,
                "Com o desafio de estruturar uma governança robusta para múltiplas unidades de negócio, "
                "a proposta era clara: garantir isolamento, segurança e rastreabilidade de custos, "
                "com uma arquitetura escalável baseada em Landing Zone na AWS.")
        
        # Shape 26 - "Você sabia?"
        elif text == "Você sabia?":
            pass  # mantém
        
        # Shape 27 - Texto do "Você sabia?"
        elif text.startswith("O monitoramento em tempo real"):
            replace_shape_text(shape,
                "Segundo a AWS, empresas que implementam Landing Zones reduzem em até 50% o tempo de "
                "provisionamento de novas contas e ambientes, além de diminuir significativamente o risco "
                "de incidentes de segurança por configurações incorretas. A governança centralizada também "
                "pode gerar economia de até 30% em custos operacionais de cloud.")
    
    # Substituir logo do cliente (shape 13 no slide 1)
    logo_path = "public/content-images/cases/grupo-autoglass-eleva-governanca-e-seguranca-na-aws-com-implementacao-de-landing-zone/logo-autoglass-RGB_LOGO-ORIGINAL-para-fundos-claros-3.png"
    if os.path.exists(logo_path):
        replace_image(slide1, 13, logo_path)
        print(f"  ✅ Logo substituído: {logo_path}")
    else:
        print(f"  ⚠️ Logo não encontrado: {logo_path}")
    
    # ============================================================
    # SLIDE 2
    # ============================================================
    slide2 = prs.slides[1]
    
    for shape in slide2.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        
        text = shape.text.strip()
        
        # Shape 12 - Título "ARQUITETURA FLUIDA..."
        if text.startswith("ARQUITETURA FLUIDA"):
            replace_shape_text(shape,
                "LANDING ZONE: ESTRUTURA SEGURA E ESCALÁVEL")
        
        # Shape 10 - Texto descritivo da arquitetura
        elif text.startswith("Tudo  começa"):
            replace_shape_text(shape,
                "A solução incluiu a criação de Unidades Organizacionais (OUs) segmentadas: uma OU exclusiva "
                "para a conta Master, outra para serviços de infraestrutura e uma para contas excluídas, "
                "com políticas de segurança específicas para cada uma.\n\n"
                "Foram criadas contas dedicadas para gerenciar backups, acessos via IAM Identity Center e "
                "configurações de rede. A configuração de serviços nativos como AWS CloudTrail, AWS Config, "
                "AWS GuardDuty, AWS Security Hub e AWS Backup elevou a maturidade de segurança do ambiente. "
                "Políticas de backups e tags foram implementadas para centralizar cópias de segurança e "
                "rastrear custos por BU.")
        
        # Shape 11 - Título + texto dos resultados
        elif text.startswith("RESULTADOS NA PRÁTICA"):
            replace_multiformat_text(shape, [
                ("RESULTADOS NA PRÁTICA: DA COMPLEXIDADE À GOVERNANÇA TOTAL", True),
                ("Com a nova Landing Zone, o Grupo Autoglass passou a operar com BUs em contas AWS "
                 "exclusivas, garantindo isolamento e segurança entre os recursos. A aplicação de SCPs "
                 "e a configuração do AWS Control Tower e Organizations aumentaram a governança do ambiente.\n\n"
                 "A configuração de GuardDuty, Security Hub e CloudTrail elevou os níveis de segurança e "
                 "monitoramento. A política de tags permitiu alocação precisa dos custos por BU, facilitando "
                 "a gestão financeira. A funcionalidade Account Factory do Control Tower agilizou a criação "
                 "de novas contas com consistência e segurança.\n"
                 "Mais governança, mais autonomia, mais segurança", False),
            ])
    
    # ============================================================
    # SLIDE 3
    # ============================================================
    slide3 = prs.slides[2]
    
    for shape in slide3.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        
        text = shape.text.strip()
        
        # Shape 9 - Título "PARCERIA QUE FAZ A DIFERENÇA"
        if text == "PARCERIA QUE FAZ A DIFERENÇA":
            pass  # mantém
        
        # Shape 13 - Depoimento (não tem depoimento para Autoglass, usar texto da conclusão)
        elif text.startswith('"A DataRain tem sido'):
            replace_multiformat_text(shape, [
                ('"Um dos grandes diferenciais desse projeto foi o atendimento próximo e colaborativo '
                 'da equipe da dataRain. Em vez de simplesmente substituir o time do cliente, trabalharam '
                 'lado a lado com nossos profissionais, capacitando-os e transferindo conhecimento ao longo '
                 'de todo o processo."', False),
                ("\n- Grupo Autoglass", True),
            ])
        
        # Shape 8 - Título + texto "AUTONOMIA E ESCALABILIDADE"
        elif text.startswith("AUTONOMIA E ESCALABILIDADE"):
            replace_multiformat_text(shape, [
                ("AUTONOMIA E ESCALABILIDADE", True),
                ("O trabalho colaborativo e a transferência de conhecimento garantiram que o time do "
                 "Grupo Autoglass estivesse preparado para gerenciar e operar a nova infraestrutura "
                 "com autonomia.\n\n"
                 "Hoje, a organização opera com um nível de segurança acima da média da maioria das "
                 "empresas, com um time maduro e capacitado, capaz de administrar e evoluir o ambiente "
                 "com excelência, reforçando a posição de liderança no setor automotivo.", False),
            ])
    
    prs.save(output_path)
    print(f"✅ Arquivo gerado: {output_path}")


if __name__ == "__main__":
    generate_autoglass_case()
